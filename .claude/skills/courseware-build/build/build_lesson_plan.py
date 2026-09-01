#!/usr/bin/env python3
"""Generate the WSQ AI Security Governance for Businesses Lesson Plan (LP) DOCX.

Cover page + Document Version Control Record + auto TOC + Arial 11pt body +
colour-coded 2-day schedule tables (9:30am-6:30pm, 8 training hours/day, 1h lunch,
tea within). Topics/labs and the SLIDE NUMBERS come from course_data + the domain
data files + the built deck, so the LP stays aligned with the deck, guide and labs.
"""
import os, sys, json, re
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
import course_data as C
from data_domain1 import DOMAIN1; from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3; from data_domain4 import DOMAIN4
ACT=DOMAIN1+DOMAIN2+DOMAIN3+DOMAIN4
import prodoc
def _find_repo(start):
    env=os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env): return env
    d=start
    for _ in range(8):
        d=os.path.dirname(d)
        if os.path.isdir(os.path.join(d,"courseware")) and os.path.isdir(os.path.join(d,"labs")): return d
    return os.path.dirname(os.path.dirname(HERE))
REPO=_find_repo(HERE); ASSETS=os.path.join(os.path.dirname(HERE),"assets")

def _asset(name):
    """Resolve a brand asset: the course's own courseware/assets first, then the skill's."""
    for base in (os.path.join(REPO,"courseware","assets"), ASSETS):
        q=os.path.join(base,name)
        if os.path.exists(q): return q
    return None


BRAND=RGBColor(0x1F,0x6F,0xEB); DARK=RGBColor(0x11,0x18,0x27); GREY=RGBColor(0x55,0x5B,0x66)
HEADER_FILL="1F6FEB"; TOPIC_FILL="E8F0FE"; BREAK_FILL="FFF4E5"; LUNCH_FILL="FDE9D9"; ASSESS_FILL="E8F7EE"

# ---- slide numbers read live from the built deck (never hand-maintained) ----
def slide_map():
    deck=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
    labs={}; secs={}
    try:
        from pptx import Presentation
        p=Presentation(deck)
        for i,s in enumerate(p.slides,1):
            j=" ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for m in re.finditer(r"LAB (\d+)", j):
                n=int(m.group(1)); labs.setdefault(n,i)
            for m in re.finditer(r"TOPIC (\d\d)\b", j):
                n=int(m.group(1)); secs.setdefault(n,i)
    except Exception as e:
        print("  ! could not read deck for slide numbers:",e)
    return labs,secs
LABS,SECS=slide_map()

def lab_titles(nums):
    out=[]
    for a in ACT:
        if a["num"] in nums:
            sl=LABS.get(a["num"])
            out.append(f"Lab {a['num']}: {a['title']}"+(f" (slide {sl})" if sl else ""))
    return "; ".join(out)
def topic_ref(n):
    s=SECS.get(n)
    return f" (slides {s}+)" if s else ""

# ------------------------------------------------ schedule: 2 days x 480 training minutes
SCHEDULE = {
 1: (C.DAY_THEMES[1], [
    ("9:30","10:00",30,"admin","Welcome, course introduction, learning outcomes, ground rules and mandatory digital attendance (AM)"),
    ("10:00","11:15",75,"topic","Core concepts: why AI changes the security picture; the expanded attack surface; traditional vs AI-era threats; the 12 NIST GenAI risks; the lethal trifecta"+topic_ref(1)),
    ("11:15","11:30",15,"break","Tea break"),
    ("11:30","13:00",90,"lab","Topic 1 — AI Security Governance Foundations and Business Risks (concepts + demo). Hands-on: "+lab_titles([1,2])),
    ("13:00","14:00",60,"lunch","Lunch break — mandatory digital attendance (PM) on return"),
    ("14:00","15:30",90,"lab","Hands-on: "+lab_titles([3,4])),
    ("15:30","15:45",15,"break","Tea break"),
    ("15:45","18:15",150,"lab","Topic 2 — Building an AI Governance Framework: NIST AI RMF, MGF for GenAI, policy architecture, roles and RACI (concepts)"+topic_ref(2)+". Hands-on: "+lab_titles([5,6,7,8])),
    ("18:15","18:30",15,"recap","Day 1 recap, Q&A and PM digital attendance"),
 ]),
 2: (C.DAY_THEMES[2], [
    ("9:30","9:45",15,"recap","Day 1 recap and mandatory digital attendance (AM)"),
    ("9:45","11:15",90,"topic","Topic 3 — Governance Controls Across the AI Lifecycle: data governance, secure development, testing and red-teaming, deployment gates, monitoring, incident response, decommissioning (concepts + demo)"+topic_ref(3)),
    ("11:15","11:30",15,"break","Tea break"),
    ("11:30","13:00",90,"lab","Hands-on: "+lab_titles([9,10])),
    ("13:00","14:00",60,"lunch","Lunch break — mandatory digital attendance (PM) on return"),
    ("14:00","15:15",75,"lab","Hands-on: "+lab_titles([11,12])),
    ("15:15","15:30",15,"break","Tea break"),
    ("15:30","17:15",105,"lab","Topic 4 — Security Controls for Agentic AI, Risk Assessment and Implementation Roadmap: agentic threat model, identity and least agency, guardrails, permission ladders, CSA agentic RMF profile (concepts)"+topic_ref(4)+". Hands-on: "+lab_titles([13,14,15,16])),
    ("17:15","18:30",75,"lab","Course recap and consolidation of the capstone roadmap; Q&A; TRAQOM survey; Briefing for Assessment; Assessment digital attendance"),
 ]),
}

# ------------------------------------------------ build document
doc=Document()
normal=doc.styles["Normal"]; normal.font.name="Arial"; normal.font.size=Pt(11)
prodoc.style_headings(doc)

prodoc.add_cover_page(doc,"LESSON PLAN",C.TITLE,C.VERSION.lstrip("v"),
                      org_logo=_asset("tertiary-infotech-logo.png"),
                      course_logo=_asset("wsq-badge.png"), course_code=C.COURSE_CODE)
prodoc.add_version_control(doc,[
 ("1.0",C.VERSION_DATE,
  "Initial release — 2-day WSQ lesson plan for AI Security Governance for Businesses, aligned to the four topics and the 16 hands-on labs, with slide references to deck "+C.VERSION+".",
  C.TRAINER),
])
prodoc.add_toc(doc)

def H(text,level=1):
    return doc.add_heading(text,level=level)

H("Course Information",1)
info=[("Course Title",C.TITLE),("WSQ Course Reference",C.COURSE_CODE),
      ("TSC Title / Code",f"{C.TSC_TITLE} · {C.TSC_CODE} · {C.TSC_LEVEL}"),
      ("Training Provider",C.ORG+"  ("+C.UEN.replace('UEN: ','UEN ')+")"),
      ("Duration","2 days · 8 training hours per day (16 hours)"),
      ("Daily Timing","9:30 am – 6:30 pm (1-hour lunch; tea breaks within training time)"),
      ("Mode","Instructor-led classroom, synchronous Zoom or corporate on-site; hands-on governance labs using browser-based security tools"),
      ("Trainer",C.TRAINER),
      ("Trainer Slides",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")]
t=doc.add_table(rows=0,cols=2); t.style="Table Grid"
for k,v in info:
    c=t.add_row().cells; c[0].text=""; r=c[0].paragraphs[0].add_run(k); r.bold=True; r.font.size=Pt(10)
    prodoc._shade_cell(c[0],TOPIC_FILL)
    c[1].text=""; c[1].paragraphs[0].add_run(v).font.size=Pt(10)

H("Learning Outcomes",1)
doc.add_paragraph("On completion of this course, learners will be able to:")
for lo in C.LEARNING_OUTCOMES:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(lo).font.size=Pt(10.5)

H("Assessment",1)
for a in [C.ASSESSMENT["written"],C.ASSESSMENT["practical"],
          "Format: Open Book — course slides, Learner Guide, the learner's own lab outputs and approved materials only.",
          "Final assessment is administered on Day 2 after the 8 instructional hours: the Written Assessment (1 hour) followed by the Case Study (1 hour) — 120 minutes in total.",
          "Grading: Competent (C) / Not Yet Competent (NYC) on each instrument.",
          C.ASSESSMENT["note"]]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(a).font.size=Pt(10.5)

def set_cell(cell,text,bold=False,size=9.5,color=None,fill=None,align=None):
    cell.text=""; p=cell.paragraphs[0]
    if align: p.alignment=align
    r=p.add_run(text); r.bold=bold; r.font.size=Pt(size); r.font.name="Arial"
    if color: r.font.color.rgb=color
    if fill: prodoc._shade_cell(cell,fill)

KIND_FILL={"topic":TOPIC_FILL,"break":BREAK_FILL,"lunch":LUNCH_FILL,"assess":ASSESS_FILL,
           "admin":"F3F5F8","recap":"F3F5F8","lab":None}

H("Course Schedule",1)
for day,(theme,rows) in SCHEDULE.items():
    H(f"Day {day} — {theme}",2)
    tbl=doc.add_table(rows=0,cols=3); tbl.style="Table Grid"; tbl.alignment=WD_TABLE_ALIGNMENT.CENTER
    hdr=tbl.add_row().cells
    for i,htext in enumerate(["Time","Duration","Topic / Activity (with slide references)"]):
        set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
    training=0
    for start,end,mins,kind,text in rows:
        cells=tbl.add_row().cells; fill=KIND_FILL.get(kind)
        set_cell(cells[0],f"{start}–{end}",bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        set_cell(cells[1],f"{mins} min",size=9.5,fill=fill)
        set_cell(cells[2],text,bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        if kind!="lunch": training+=mins
    for row in tbl.rows:
        row.cells[0].width=Inches(1.15); row.cells[1].width=Inches(0.9); row.cells[2].width=Inches(4.75)
    p=doc.add_paragraph(); r=p.add_run(f"Total training time: {training} minutes ({training//60} hours)."); r.italic=True; r.font.size=Pt(9.5); r.font.color.rgb=GREY
    assert training==480, f"Day {day} training minutes = {training}, expected 480"
    if day==C.DAYS:
        H("Assessment Administration (Day 2, after the instructional hours)",3)
        atbl=doc.add_table(rows=0,cols=3); atbl.style="Table Grid"; atbl.alignment=WD_TABLE_ALIGNMENT.CENTER
        ah=atbl.add_row().cells
        for i,htext in enumerate(["Time","Duration","Assessment Instrument"]):
            set_cell(ah[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
        for st,en,mn,tx in [
            ("18:30","19:30",60,"Written Assessment (WA) — Short-Answer Questions (SAQ), 1 hour, open book"),
            ("19:30","20:30",60,"Case Study (CS) — one coherent governance case study, 1 hour, open book"),
        ]:
            cc=atbl.add_row().cells
            set_cell(cc[0],f"{st}–{en}",bold=True,size=9.5,fill=ASSESS_FILL)
            set_cell(cc[1],f"{mn} min",size=9.5,fill=ASSESS_FILL)
            set_cell(cc[2],tx,bold=True,size=9.5,fill=ASSESS_FILL)
        for row in atbl.rows:
            row.cells[0].width=Inches(1.15); row.cells[1].width=Inches(0.9); row.cells[2].width=Inches(4.75)
        p=doc.add_paragraph()
        r=p.add_run("Total assessment time: 120 minutes (2 hours), additional to the 8 instructional hours above. "
                    "Both instruments are open book. The Assessment digital attendance is taken before the papers "
                    "are issued. Where the schedule requires, the assessment may be administered on a separate "
                    "session; the instruments and their durations do not change.")
        r.italic=True; r.font.size=Pt(9.5); r.font.color.rgb=GREY

H("Lab Reference (aligned to topics and slide numbers)",1)
tt=doc.add_table(rows=0,cols=4); tt.style="Table Grid"
hdr=tt.add_row().cells
for i,htext in enumerate(["Topic","Weighting","Labs","Deck slides"]):
    set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for tp in C.TOPICS:
    acts=[a for a in ACT if a["topic"]==tp["num"]]
    cells=tt.add_row().cells
    set_cell(cells[0],f"Topic {tp['code']}: {tp['title']}",bold=True,size=9.5,fill=TOPIC_FILL)
    set_cell(cells[1],tp["weighting"],size=9.5,fill=TOPIC_FILL)
    set_cell(cells[2],", ".join(f"Lab {a['num']}" for a in acts),size=9.5)
    sl=[LABS.get(a["num"]) for a in acts if LABS.get(a["num"])]
    set_cell(cells[3],(f"{SECS.get(tp['num'],'')} · labs {min(sl)}–{max(sl)}" if sl else ""),size=9.5)

H("Lab Detail",1)
lt=doc.add_table(rows=0,cols=4); lt.style="Table Grid"
hdr=lt.add_row().cells
for i,htext in enumerate(["Lab","Title","Learners produce","Slide"]):
    set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for a in ACT:
    cells=lt.add_row().cells
    set_cell(cells[0],str(a["num"]),bold=True,size=9,fill=TOPIC_FILL)
    set_cell(cells[1],a["title"],size=9)
    set_cell(cells[2],a["build"],size=8.5)
    set_cell(cells[3],str(LABS.get(a["num"],"")),size=9)
for row in lt.rows:
    row.cells[0].width=Inches(0.5); row.cells[1].width=Inches(2.1)
    row.cells[2].width=Inches(3.6); row.cells[3].width=Inches(0.6)

H("Training Resources",1)
for r_ in ["Trainer Slides (PPT) and Learner Slides (PDF) — "+f"{C.SHORT_TITLE}-{C.VERSION}",
           "Learner Guide (DOCX/PDF) with the full step-by-step procedure for all 16 labs",
           "Lab pack — one folder per lab with mock data and a printable instruction PDF",
           "Browser-based lab tools (no installation required):"]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(r_).font.size=Pt(10.5)
for name,url in C.TOOLS:
    p=doc.add_paragraph(style="List Bullet 2"); p.add_run(f"{name} — {url}").font.size=Pt(10)
for r_ in ["Whiteboard/flipchart for the incident-response tabletop exercise (Lab 12)",
           "LMS/TMS portal for course materials and assessment submission — https://lms-tms.tertiaryinfotech.com"]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(r_).font.size=Pt(10.5)

prodoc.add_page_numbers(doc)
prodoc.enable_update_fields(doc)
OUT=os.path.join(REPO,"courseware",f"LP-{C.SHORT_TITLE}.docx")
doc.save(OUT)
print("Saved",OUT)
