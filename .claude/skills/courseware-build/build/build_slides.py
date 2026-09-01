#!/usr/bin/env python3
"""Generate the WSQ AI Security Governance for Businesses slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys, copy, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

# ---------------- motion: restrained transitions + build animation ----------------
# House rule: ONE transition family for the whole deck (morph is unavailable in the
# OOXML PowerPoint 2010 transition set that python-pptx can emit, so we use a short
# push/fade pair). Content slides fade; section dividers push. Nothing else moves.
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _transition(s, kind="fade", speed="med"):
    """Attach a restrained slide transition. kind: fade | push | wipe."""
    sld = s._element
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", speed)
    tr.set("advClick", "1")
    if kind == "fade":
        etree.SubElement(tr, qn("p:fade"))
    elif kind == "push":
        el = etree.SubElement(tr, qn("p:push")); el.set("dir", "l")
    elif kind == "wipe":
        el = etree.SubElement(tr, qn("p:wipe")); el.set("dir", "r")
    # keep the transition last in the slide element (schema order)
    sld.append(tr)
    return tr

# Appear-on-click build for a list of shape ids — used ONLY on process maps so the
# trainer can reveal one stage at a time. No spins, no flying, no sound.
_TIMING = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
<p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">
<p:childTnLst>{pars}</p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""

_PAR = """<p:par><p:cTn id="{i0}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
<p:childTnLst><p:par><p:cTn id="{i2}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" nodeType="{nt}">
<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
<p:set><p:cBhvr><p:cTn id="{i3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{i4}" dur="400"/>
<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>
</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>"""

def _build_on_click(s, spids):
    """Fade each shape id in on click, in order. Used sparingly (process maps)."""
    if not spids:
        return
    pars, nid = [], 10
    for k, spid in enumerate(spids):
        pars.append(_PAR.format(i0=nid, i1=nid+1, i2=nid+2, i3=nid+3, i4=nid+4,
                                spid=spid, nt="clickEffect" if k == 0 else "afterEffect"))
        nid += 10
    sld = s._element
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)
    sld.append(etree.fromstring(_TIMING.format(pars="".join(pars))))

def connector(s, x1, y1, x2, y2, color, width=Pt(2.0), arrow=True):
    """A REAL PowerPoint connector line (not a typed arrow glyph)."""
    cx, cy = min(x1, x2), min(y1, y2)
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return cn

def chevron(s, x, y, w, h, color):
    """A staged chevron shape — the real CHEVRON autoshape, not a text glyph."""
    sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def roundrect(s, x, y, w, h, color, line=None, adj=0.10):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = adj
    except Exception: pass
    return sp

def diamond(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp

def label_in(sp, text, size, color, bold=True):
    """Put centred text inside an autoshape."""
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Arial"
    return sp

PAGE={"n":1}   # the cover is slide 1 and carries no number, so numbering starts at 2
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def _ellipsis(text,limit):
    """Truncate on a WORD boundary with an ellipsis — never slice mid-word."""
    t=" ".join(str(text).split())
    if len(t)<=limit: return t
    cut=t[:limit]
    sp=cut.rfind(" ")
    if sp>limit*0.55: cut=cut[:sp]
    return cut.rstrip(" ,.;:-—(") + "…"

def _short_cmd(cmd,limit=30):
    """Shorten a shell command for a caption. Paths/URLs have no spaces, so a
    word-boundary ellipsis cuts them mid-token — drop the runner prefix and keep the
    meaningful tail (the script/target) instead."""
    c=" ".join(str(cmd).split())
    for pre in ("uv run python ","uv run ","python3 ","python ","bash ","sh "):
        if c.startswith(pre): c=c[len(pre):]; break
    if len(c)<=limit: return c
    parts=c.split(" ")
    head=parts[0]
    rest=" ".join(parts[1:]).strip()
    if rest:
        # a URL argument → keep the verb and the host, drop the path
        m=re.match(r'https?://([^/\s]+)',rest)
        if m:
            for cand in (f"{head} {m.group(1)}/…", f"{head} {m.group(1)}"):
                if len(cand)<=limit: return cand
        # "git clone …/repo" — keep the verb plus the final path segment
        tail=rest.rstrip("/").split("/")[-1]
        cand=f"{head} …/{tail}" if "/" in rest else f"{head} {tail}"
        if len(cand)<=limit: return cand
        if len(head)+2<=limit: return _ellipsis(head,limit-2)+" …"
    if "/" in c:                   # a bare path → keep the last segment
        tail=c.rstrip("/").split("/")[-1]
        if len(tail)+2<=limit: return "…/"+tail
    return _ellipsis(c,limit)

def _fit_title(title,size=29):
    """Shrink long titles so they never wrap into the hairline rule below."""
    n=len(title)
    if n<=52: return size
    if n<=66: return 25
    if n<=82: return 22
    return 20

def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.88),Inches(11.9),Inches(0.78),
        [[(title,_fit_title(title),INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    for base in (os.path.join(REPO,"courseware","assets"), ASSETS):
        p=os.path.join(base,name)
        if os.path.exists(p): return p
    return None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    wsq=_logo("wsq-badge.png")   # second cover logo: the WSQ / SkillsFuture credential mark
    if wsq: s.shapes.add_picture(wsq,Inches(2.35),Inches(0.78),height=Inches(0.9))
    # course badge (top-right) — Gemini Agent ADK badge, else text fallback
    badge=_logo("ai-security-governance-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.35),Inches(0.6),width=Inches(2.2))
    else:
        rect(s,Inches(10.62),Inches(0.62),Inches(2.0),Inches(1.2),BLUE)
        txt(s,Inches(10.62),Inches(0.76),Inches(2.0),Inches(0.5),[[("AI SECURITY",17,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(10.62),Inches(1.28),Inches(2.0),Inches(0.42),[[("GOVERNANCE",11,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    # The title band is 1.55in tall before the subtitle starts. Shrink a long title
    # until it fits that band, so a 3-line title can never collide with the subtitle.
    TW=Inches(11.4); tsz=40
    while tsz>22 and _wrapped_lines(title,TW,tsz)*(tsz*1.20) > (Inches(1.5)/12700.0):
        tsz-=2
    txt(s,Inches(1.25),Inches(3.0),TW,Inches(1.55),[[(title,tsz,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.62),Inches(11),Inches(0.9),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]

def _wrapped_lines(text,width_emu,pt):
    """Estimate how many lines `text` wraps to in a box `width_emu` wide at `pt`.
    Arial averages about 0.50 x font-size per character, so chars-per-line is
    (width in points) / (0.50 * pt). Used to auto-fit tile text instead of
    letting a 4th line spill past the card edge."""
    import math as _m
    wpt=width_emu/12700.0                      # EMU -> points
    cpl=max(int(wpt/(0.55*pt)),8)
    n=0
    for para in str(text).split("\n"):
        n+=max(1,_m.ceil(len(para)/cpl))
    return n

def _fit_tile(title_txt,body_txt,tw,ch,tsize,bsize,pad_emu=None):
    """Shrink the tile's title/body until the wrapped text fits the TEXTBOX height.
    `ch` is the card height; the textbox inside it is shorter (it is inset), and the
    text frame adds its own top/bottom insets — so measure against the real box, not
    the card, or the last line spills past the card edge.
    Returns the (title_pt, body_pt) that fit."""
    inner=ch-(pad_emu if pad_emu is not None else Inches(0.2))
    avail=(inner/12700.0)-8                      # card height in points, less padding
    t,b=tsize,bsize
    for _ in range(9):
        lines_t=_wrapped_lines(title_txt,tw,t) if title_txt else 0
        lines_b=_wrapped_lines(body_txt,tw,b) if body_txt else 0
        need=lines_t*(t*1.22)+lines_b*(b*1.22)+(3 if title_txt and body_txt else 0)
        if need<=avail or b<=8: break
        b-=1
        if t>b+2: t-=1
    return t,b

def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            ts,bs=_fit_tile(it[0],it[1],tw,ch,size+2,size-2)
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],ts,INK,True)],[(it[1],bs,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            _,bs=_fit_tile("",it,tw,ch,size,size)
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,bs,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def process_map(title,stages,kicker=None,color=BLUE,synthesis=None,animate=True):
    """STAGED PROCESS MAP — real rounded-rect stages joined by real connectors with
    arrowheads, each stage numbered, optional synthesis band. stages: list of
    (label, detail). This replaces flow_h wherever a genuine process is taught."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(stages); X0=Inches(0.85); TOTW=Inches(11.63)
    gap=Inches(0.42); cw=int((TOTW-gap*(n-1))/n)
    y=Inches(2.35); ch=Inches(2.35) if synthesis else Inches(3.0)
    spids=[]
    for i,st in enumerate(stages):
        lbl,detail=(st if isinstance(st,tuple) else (st,""))
        x=int(X0+(cw+gap)*i)
        box=roundrect(s,x,y,cw,ch,LIGHT,line=LINE)
        rect(s,x,y,cw,Inches(0.11),color)
        bd=Inches(0.62)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.34)),bd,bd,[[(str(i+1),24,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        # label sits in its own fixed band; the caption gets a single line beneath it
        lbl=" ".join(str(lbl).split())
        # Label wraps inside its band; shrink to fit rather than clipping the words off.
        _lw=cw-Inches(0.24)
        _,lsz=_fit_tile("",lbl,_lw,Inches(1.06),13,13,pad_emu=0)
        txt(s,x+Inches(0.12),int(y+Inches(1.00)),_lw,Inches(1.06),
            [[(lbl,lsz,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        if detail:
            # ONE line only, sized to the card — a 2-line caption overflows the card bottom
            det=" ".join(str(detail).split())
            _dw2=cw-Inches(0.16)
            _,dsz=_fit_tile("",det,_dw2,Inches(0.44),9,9,pad_emu=0)
            txt(s,x+Inches(0.08),int(y+ch-Inches(0.56)),_dw2,Inches(0.44),
                [[(det,dsz,GREY,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=0)
        spids.append(box.shape_id)
        if i<n-1:   # REAL connector in the gap, with an arrowhead
            cy=int(y+ch/2)
            cn=connector(s,int(x+cw+Inches(0.06)),cy,int(x+cw+gap-Inches(0.06)),cy,color)
            spids.append(cn.shape_id)
    if synthesis:
        by=int(y+ch+Inches(0.34))
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.15),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.15),color)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.32),
            [[(synthesis[0].upper(),11,color,True)]])
        _,_ssz=_fit_tile("",synthesis[1],Inches(11.1),Inches(0.62),13,13,pad_emu=0)
        txt(s,Inches(1.15),int(by+Inches(0.46)),Inches(11.1),Inches(0.62),
            [[(synthesis[1],_ssz,INK,False)]])
    if animate: _build_on_click(s,spids)
    footer(s); return s

def decision_map(title,question,yes,no,kicker=None,color=VIOLET,note=None):
    """A real decision diamond with two branches drawn as connectors — used for
    'which pattern do I choose' teaching moments."""
    s=head(slide(),title,kicker,kcolor=color)
    # a diamond's usable text area is ~50% of its box — size generously or text spills
    # past the facets. 4.6 x 2.7 fits 2-3 short lines at 12pt.
    dx,dy,dw,dh=Inches(0.85),Inches(2.75),Inches(4.6),Inches(2.7)
    d=diamond(s,dx,dy,dw,dh,color); label_in(d,question,12,WHITE)
    bx=Inches(6.1); bw=Inches(6.35); bh=Inches(1.5)
    ys=[Inches(2.15),Inches(4.35)]
    for (hdr,items),by,col in zip([yes,no],ys,[TEAL,AMBER]):
        b=roundrect(s,bx,by,bw,bh,LIGHT,line=LINE)
        rect(s,bx,by,Inches(0.11),bh,col)
        txt(s,bx+Inches(0.3),int(by+Inches(0.16)),bw-Inches(0.55),Inches(0.4),[[(hdr,15,col,True)]])
        txt(s,bx+Inches(0.3),int(by+Inches(0.6)),bw-Inches(0.55),Inches(0.82),[[(items,12,INK,False)]])
        connector(s,int(dx+dw),int(dy+dh/2),bx,int(by+bh/2),col)
    if note:
        rect(s,Inches(0.85),Inches(6.15),Inches(11.63),Inches(0.72),LIGHT)
        txt(s,Inches(1.15),Inches(6.28),Inches(11.1),Inches(0.5),[[(note,12,GREY,False)]])
    footer(s); return s

def compare_table(title,headers,rows,kicker=None,accent=BLUE,note=None):
    """A real comparison matrix — the substantive alternative to two bullet columns."""
    s=head(slide(),title,kicker,kcolor=accent)
    X0=Inches(0.85); TOTW=Inches(11.63); ncol=len(headers)
    first=int(TOTW*0.26); rest=int((TOTW-first)/(ncol-1))
    widths=[first]+[rest]*(ncol-1)
    y=Inches(1.95); hh=Inches(0.52)
    x=X0
    for i,h in enumerate(headers):
        col=accent if i==0 else PALETTE[(i-1)%len(PALETTE)]
        rect(s,x,y,widths[i],hh,col)
        txt(s,x+Inches(0.14),y,widths[i]-Inches(0.28),hh,[[(h,13,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
        x+=widths[i]
    # body must end above the note band AND the footer (7.05). Budget it explicitly.
    ry=int(y+hh)
    BOTTOM=Inches(6.88)-(Inches(1.24) if note else Emu(0))
    avail=BOTTOM-ry
    rh=int(min(Inches(0.78),avail/max(len(rows),1)))
    for r,row in enumerate(rows):
        x=X0
        for i,cell in enumerate(row):
            fill=LIGHT if r%2==0 else WHITE
            rect(s,x,ry,widths[i],rh,fill,line=LINE)
            bold=(i==0)
            _cw=widths[i]-Inches(0.28)
            _,_csz=_fit_tile("",cell,_cw,rh,11.5,11.5,pad_emu=0)
            txt(s,x+Inches(0.14),ry,_cw,rh,
                [[(cell,_csz,INK if bold else GREY,bold)]],anchor=MSO_ANCHOR.MIDDLE)
            x+=widths[i]
        ry+=rh
    if note:
        rect(s,X0,int(ry+Inches(0.22)),TOTW,Inches(0.92),LIGHT)
        rect(s,X0,int(ry+Inches(0.22)),Inches(0.11),Inches(0.92),accent)
        txt(s,X0+Inches(0.3),int(ry+Inches(0.34)),TOTW-Inches(0.6),Inches(0.66),
            [[("WHEN IT MATTERS  ",11,accent,True),(note,12,INK,False)]])
    footer(s); return s

def worked_example(title,intro,code,explain,kicker=None,accent=TEAL):
    """A worked example: the code on the left, the line-by-line reading on the right.
    This is what turns a decorative lab slide into a teaching slide."""
    s=head(slide(),title,kicker,kcolor=accent)
    txt(s,Inches(0.85),Inches(1.9),Inches(11.63),Inches(0.46),[[(intro,15,GREY,False)]])
    cx,cw=Inches(0.85),Inches(6.5)
    rect(s,cx,Inches(2.5),cw,Inches(4.15),RGBColor(0x0B,0x12,0x20))
    tb=s.shapes.add_textbox(cx+Inches(0.22),Inches(2.62),cw-Inches(0.44),Inches(3.9))
    tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(code):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(2)
        r=p.add_run(); r.text=ln; r.font.size=Pt(11); r.font.name="Consolas"
        col=RGBColor(0x9C,0xDC,0xFE)
        st=ln.strip()
        if st.startswith("#"): col=RGBColor(0x6A,0x99,0x55)
        elif "=" in ln and not st.startswith(("def","class")): col=RGBColor(0xD4,0xD4,0xD4)
        if st.startswith(("def ","class ","from ","import ")): col=RGBColor(0xC5,0x86,0xC0)
        r.font.color.rgb=col
    ex,ew=Inches(7.65),Inches(4.83)
    for i,(lbl,body) in enumerate(explain[:4]):
        y=int(Inches(2.5)+(Inches(1.0)+Inches(0.05))*i)
        col=PALETTE[i%len(PALETTE)]
        rect(s,ex,y,ew,Inches(1.0),LIGHT); rect(s,ex,y,Inches(0.09),Inches(1.0),col)
        txt(s,ex+Inches(0.26),int(y+Inches(0.1)),ew-Inches(0.45),Inches(0.32),[[(lbl,12,col,True)]])
        txt(s,ex+Inches(0.26),int(y+Inches(0.42)),ew-Inches(0.45),Inches(0.52),[[(body,11,INK,False)]])
    footer(s); return s

def steps_slide(act_title,steps,kicker,accent=TEAL,part=None,start=1):
    """Substantive lab procedure: numbered steps WITH their commands, up to 5 per
    slide. Replaces the old one-sentence step slide."""
    s=head(slide(),act_title+(f" — {part}" if part else ""),kicker,kcolor=accent)
    y0=Inches(1.92); n=len(steps); gapy=Inches(0.1)
    AVAIL=Inches(4.92)          # 1.92 → 6.84, clear of the 7.05 footer
    rh=int(min(Inches(1.12),(AVAIL-gapy*(n-1))/max(n,1)))
    for i,(text,cmd) in enumerate(steps):
        y=int(y0+(rh+gapy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,Inches(0.85),y,Inches(11.63),rh,LIGHT); rect(s,Inches(0.85),y,Inches(0.09),rh,col)
        bd=Inches(0.4)
        oval(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,col)
        txt(s,Inches(1.06),int(y+rh/2-bd/2),bd,bd,[[(str(start+i),13,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        if cmd:
            txt(s,Inches(1.62),int(y+Inches(0.08)),Inches(10.6),Inches(0.34),[[(text,12.5,INK,True)]])
            rect(s,Inches(1.62),int(y+Inches(0.44)),Inches(10.5),int(rh-Inches(0.54)),RGBColor(0x0B,0x12,0x20))
            one=cmd.split("\n")[0]
            if len(one)>96: one=one[:93]+"..."
            txt(s,Inches(1.78),int(y+Inches(0.44)),Inches(10.2),int(rh-Inches(0.54)),
                [[("$ "+one,10.5,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
        else:
            txt(s,Inches(1.62),y,Inches(10.6),rh,[[(text,12.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s

def chart_slide(title,categories,series,kicker=None,accent=BLUE,
                kind="bar",insight=None,number_format='0'):
    """NATIVE PowerPoint chart (fully editable, not a picture) + an insight band.
    kind: bar | column | line | pie | doughnut. series: list of (name, values)."""
    s=head(slide(),title,kicker,kcolor=accent)
    cd=ChartData(); cd.categories=categories
    for nm,vals in series: cd.add_series(nm,vals,number_format)
    ctype={"bar":XL_CHART_TYPE.BAR_CLUSTERED,"column":XL_CHART_TYPE.COLUMN_CLUSTERED,
           "line":XL_CHART_TYPE.LINE_MARKERS,"pie":XL_CHART_TYPE.PIE,
           "doughnut":XL_CHART_TYPE.DOUGHNUT}.get(kind,XL_CHART_TYPE.COLUMN_CLUSTERED)
    ch_h=Inches(3.55) if insight else Inches(4.75)
    gf=s.shapes.add_chart(ctype,Inches(0.85),Inches(1.95),Inches(11.63),ch_h,cd)
    ch=gf.chart
    ch.has_title=False
    ch.font.size=Pt(12); ch.font.name="Arial"; ch.font.color.rgb=INK
    if kind in ("pie","doughnut") or len(series)>1:
        ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    else:
        ch.has_legend=False
    # house palette per point (pie/doughnut) or per series (bar/column/line)
    try:
        if kind in ("pie","doughnut"):
            pts=ch.plots[0]
            for i,pt in enumerate(pts.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb=PALETTE[i%len(PALETTE)]
        else:
            for i,sr in enumerate(ch.series):
                col=PALETTE[i%len(PALETTE)]
                if kind=="line":
                    sr.format.line.color.rgb=col; sr.format.line.width=Pt(2.5)
                else:
                    sr.format.fill.solid(); sr.format.fill.fore_color.rgb=col
    except Exception:
        pass
    try:
        pl=ch.plots[0]; pl.has_data_labels=True
        dl=pl.data_labels; dl.font.size=Pt(10); dl.font.color.rgb=INK
        dl.number_format=number_format; dl.number_format_is_linked=False
        if kind in ("pie","doughnut"): dl.position=XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    if insight:
        by=Inches(5.72)
        rect(s,Inches(0.85),by,Inches(11.63),Inches(1.1),LIGHT)
        rect(s,Inches(0.85),by,Inches(0.11),Inches(1.1),accent)
        txt(s,Inches(1.15),int(by+Inches(0.12)),Inches(11.1),Inches(0.3),
            [[("WHAT THE DATA SHOWS",11,accent,True)]])
        txt(s,Inches(1.15),int(by+Inches(0.44)),Inches(11.1),Inches(0.6),
            [[(insight,12,INK,False)]])
    footer(s); return s

def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker,objective=None,test=None):
    """Lab briefing — now a full teaching slide: the tag chip, the description, and a
    3-tile band covering objective / deliverable / toolchain, plus the success test."""
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),TEAL)
    txt(s,Inches(0.85),Inches(1.88),Inches(1.7),Inches(0.46),[[(tag,15,WHITE,True)]],
        align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if objective:
        txt(s,Inches(2.72),Inches(1.88),Inches(9.7),Inches(0.46),
            [[(objective,12,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    _dw=Inches(11.63)
    _,_dsz=_fit_tile("",desc,_dw,Inches(1.15),17,17,pad_emu=0)
    txt(s,Inches(0.85),Inches(2.5),_dw,Inches(1.15),[[(desc,_dsz,INK,False)]])
    tiles=[(BLUE,"YOU'LL BUILD",build),(TEAL,"TOOLCHAIN",services),
           (VIOLET,"DONE WHEN",test or "The lab runs end to end without error.")]
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for (col,lbl,body),x in zip(tiles,xs):
        rect(s,x,Inches(3.8),tw,Inches(2.05),LIGHT); rect(s,x,Inches(3.8),tw,Inches(0.1),col)
        txt(s,x+Inches(0.24),Inches(3.98),tw-Inches(0.45),Inches(0.34),[[(lbl,11,col,True)]])
        _bw=tw-Inches(0.45)
        _,_bsz=_fit_tile("",body,_bw,Inches(1.4),12,12,pad_emu=0)
        txt(s,x+Inches(0.24),Inches(4.34),_bw,Inches(1.4),[[(body,_bsz,INK,False)]])
    footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker,troubleshoot=None):
    """Verification slide — the success criterion PLUS a troubleshooting band, so it
    teaches diagnosis rather than stating one sentence."""
    s=head(slide(),act_title,kicker,TEAL)
    GREEN=RGBColor(0x12,0x7A,0x3E)
    rect(s,Inches(0.85),Inches(1.95),Inches(11.63),Inches(2.15),RGBColor(0xE8,0xF7,0xEE))
    rect(s,Inches(0.85),Inches(1.95),Inches(0.11),Inches(2.15),GREEN)
    txt(s,Inches(1.2),Inches(2.12),Inches(11),Inches(0.44),[[("✅  Expected result",17,GREEN,True)]])
    txt(s,Inches(1.2),Inches(2.62),Inches(11.0),Inches(1.35),[[(text,15,INK,False)]])
    tb=troubleshoot or [
        ("Nothing happens","Check the .env file is in the labs folder and the key has no quotes or spaces."),
        ("Auth or 401 error","Re-copy the API key from AI Studio; confirm GOOGLE_GENAI_USE_VERTEXAI=0."),
        ("ModuleNotFoundError","Run uv sync again, and prefix commands with uv run so the venv is used."),
    ]
    txt(s,Inches(0.85),Inches(4.32),Inches(11.63),Inches(0.34),
        [[("IF IT DOESN'T WORK",11,AMBER,True)]])
    tw=Inches(3.71); xs=[Inches(0.85),Inches(4.81),Inches(8.77)]
    for i,(sym,fix) in enumerate(tb[:3]):
        x=xs[i]
        rect(s,x,Inches(4.7),tw,Inches(1.72),LIGHT); rect(s,x,Inches(4.7),tw,Inches(0.09),AMBER)
        txt(s,x+Inches(0.24),Inches(4.87),tw-Inches(0.45),Inches(0.36),[[(sym,12.5,INK,True)]])
        txt(s,x+Inches(0.24),Inches(5.26),tw-Inches(0.45),Inches(1.05),[[(fix,11,GREY,False)]])
    footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1


def screenshot_slide(title,image,steps,kicker=None,accent=BLUE,note=None):
    """A REAL screenshot on the left with numbered captions on the right.
    Used wherever the house standard asks for a visual rather than a text link."""
    s=head(slide(),title,kicker,kcolor=accent)
    img=_logo(image)
    IX,IY,IW=Inches(0.85),Inches(1.95),Inches(6.7)
    BODY_H=Inches(4.18) if note else Inches(4.72)
    if img:
        from PIL import Image as _PILImage
        try:
            w,h=_PILImage.open(img).size; ih=int(IW*h/w)
        except Exception:
            ih=int(IW*0.62)
        if ih>BODY_H:
            ih=BODY_H; IW=int(BODY_H*w/h)
        rect(s,IX-Inches(0.06),IY-Inches(0.06),IW+Inches(0.12),ih+Inches(0.12),LINE)
        s.shapes.add_picture(img,IX,IY,width=IW,height=ih)
    else:
        rect(s,IX,IY,IW,Inches(4.2),LIGHT,line=LINE)
    x=Inches(7.85); w=Inches(4.63); y=Inches(1.95)
    n=len(steps); gy=Inches(0.14); ch=int((BODY_H-gy*(n-1))/max(n,1))
    for i,st in enumerate(steps):
        yy=int(y+(ch+gy)*i)
        rect(s,x,yy,w,ch,LIGHT); rect(s,x,yy,Inches(0.09),ch,accent)
        bd=Inches(0.42)
        oval(s,x+Inches(0.22),int(yy+ch/2-bd/2),bd,bd,accent)
        txt(s,x+Inches(0.22),int(yy+ch/2-bd/2),bd,bd,[[(str(i+1),14,WHITE,True)]],
            align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        _tw=w-Inches(0.95)
        _,_bs=_fit_tile("",st,_tw,int(ch-Inches(0.12)),13,13)
        txt(s,x+Inches(0.78),int(yy+Inches(0.06)),_tw,int(ch-Inches(0.12)),
            [[(st,_bs,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if note:
        ny=Inches(6.28)
        rect(s,Inches(0.85),ny,Inches(11.63),Inches(0.62),LIGHT)
        rect(s,Inches(0.85),ny,Inches(0.09),Inches(0.62),accent)
        txt(s,Inches(1.12),int(ny+Inches(0.06)),Inches(11.2),Inches(0.5),
            [[(note,11,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s

def gallery_slide(title,items,kicker=None,accent=BLUE,note=None):
    """A 2x2 grid of real screenshots, each with a caption — the lab tool gallery."""
    s=head(slide(),title,kicker,kcolor=accent)
    X0,Y0=Inches(0.85),Inches(1.95); CW,CH=Inches(5.68),Inches(2.08); GX,GY=Inches(0.27),Inches(0.2)
    for i,(img,cap,sub) in enumerate(items[:4]):
        r,c=i//2,i%2
        x=int(X0+(CW+GX)*c); y=int(Y0+(CH+GY)*r)
        rect(s,x,y,CW,CH,LIGHT,line=LINE)
        pth=_logo(img)
        iw=Inches(2.55); ih=int(CH-Inches(0.3))
        if pth:
            s.shapes.add_picture(pth,x+Inches(0.14),y+Inches(0.15),width=iw,height=ih)
        col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,Inches(0.09),CH,col)
        tx=x+iw+Inches(0.32); tw=CW-iw-Inches(0.46)
        txt(s,tx,y+Inches(0.24),tw,Inches(0.5),[[(cap,15,INK,True)]])
        _,_gs=_fit_tile("",sub,tw,int(CH-Inches(0.95)),11,11)
        txt(s,tx,y+Inches(0.78),tw,int(CH-Inches(0.95)),[[(sub,_gs,GREY,False)]])
    if note:
        ny=Inches(6.35)
        rect(s,Inches(0.85),ny,Inches(11.63),Inches(0.56),LIGHT)
        rect(s,Inches(0.85),ny,Inches(0.09),Inches(0.56),accent)
        txt(s,Inches(1.12),int(ny+Inches(0.04)),Inches(11.2),Inches(0.48),
            [[(note,10,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s

# ============================================================ BUILD
cover()


# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
tile_grid("Digital Attendance (Mandatory)",[
 ("Three times a day","Take the AM, PM and Assessment digital attendance — mandatory for every WSQ-funded course."),
 ("Trainer shows the QR","The trainer or administrator displays the digital attendance QR code from the SSG portal."),
 ("Scan and submit","Scan the QR code with your mobile phone camera and submit your attendance."),
 ("75% minimum","A minimum of 75% attendance is required to be eligible for assessment and funding.")],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=2,size=15)
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Expertise","AI governance and security, responsible AI, cybersecurity risk management and applied machine learning."),
  ("Delivers","WSQ courses on AI security governance, responsible AI, data protection and generative AI for business."),
  ("Credentials","PhD; ACTA/ACLP certified; certified in AI Ethics & Governance, Explainable AI and Safeguarding AI.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name, organisation and role.",
 "Does your organisation use AI today — and does anyone govern it?",
 "What worries you most about AI in your business: data leakage, wrong decisions, regulators, or something else?",
 "What you would like to be able to do differently after these two days."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
screenshot_slide("Download Your Course Material","lms-login.png",[
 "Open https://lms-tms.tertiaryinfotech.com in your browser.",
 "Sign in with the e-mail you used to register — request an OTP, or use your password.",
 "Open 'WSQ - AI Security Governance for Businesses' from your course list.",
 "Download the Trainer Slides, Learner Guide and Lesson Plan (PDF).",
 "Download the Labs folder — every lab has its own mock data and instruction PDF.",
 "Keep them open: you may use these materials during the open-book assessment."],
 kicker="LMS / TMS  ·  lms-tms.tertiaryinfotech.com",accent=BLUE,
 note="This is also where you submit your completed assessment answers, and where the TRAQOM feedback survey is found.")
tile_grid("Skills Framework Alignment",[
 ("TSC Title", C.TSC_TITLE),
 ("TSC Code", C.TSC_CODE),
 ("Proficiency Level", C.TSC_LEVEL),
 ("Knowledge","Ethical frameworks and guidelines for AI-powered tools, including legal requirements for data use."),
 ("Ability","Evaluate ethical frameworks and guidelines in AI development and deployment."),
 ("Ability","Develop and champion strategies that apply governance principles to reduce risk in deployment.")],
 kicker="SKILLS FRAMEWORK  ·  TSC",cols=2,size=14,accent=VIOLET)
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Digital Attendance (AM) · Introductions · Learning Outcomes",1),
 ("Topic 1: Foundations and Business Risks (Labs 1–4)",1),
 ("Lunch Break · Digital Attendance (PM)",1),
 ("Topic 2: Building the Governance Framework (Labs 5–8)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Digital Attendance (AM) · Day 1 recap",1),
 ("Topic 3: Governance Controls Across the AI Lifecycle (Labs 9–12)",1),
 ("Lunch Break · Digital Attendance (PM)",1),
 ("Topic 4: Agentic AI Security, Risk Assessment & Roadmap (Labs 13–16)",1),
 ("TRAQOM Survey · Digital Attendance (Assessment) · Final Assessment (WA + CS)",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Learning Outcomes",list(C.LO_SHORT),
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
tile_grid("Course Outline",[
 (f"Topic {t['code']} — {t['title']}", t["subtitle"]) for t in C.TOPICS],
 kicker="FOUR TOPICS  ·  16 HANDS-ON LABS",cols=1,size=15)
tile_grid("Briefing for Assessment",[
 ("Do · Clear your desk","Place phones and other materials under the table or on the floor."),
 ("Don't · No recording","No photos or recording of assessment scripts."),
 ("Don't · No discussion","Work individually — no discussion during the assessment."),
 ("Do · Black or blue pen","Use a black or blue pen for hard-copy assessments."),
 ("Don't · No correction fluid","No liquid paper or correction tape may be used."),
 ("Do · Stop on time","Scripts are collected when time is up.")],
 kicker="BEFORE YOU START",cols=2,size=14,accent=AMBER)
tile_grid("Assessment",[
 ("Written Assessment (WA)","Short-Answer Questions (SAQ) · 1 hour · open book. Tests the underpinning governance knowledge taught in the slides."),
 ("Case Study (CS)","One coherent governance case study · 1 hour · open book. Tests the abilities you built in the labs."),
 ("Open book","You may use the course slides, the Learner Guide, your lab outputs and approved materials only."),
 ("Eligibility","A minimum of 75% attendance is required to be eligible for assessment and funding."),
 ("Result","You are assessed as Competent (C) or Not Yet Competent (NYC) on each instrument."),
 ("Appeals","An appeal process is available if you wish to contest an assessment outcome.")],
 kicker="FINAL ASSESSMENT",cols=2,size=14)
process_map("Assessment Flow",[
 ("TRAQOM survey","Scan the QR on the LMS"),
 ("Digital attendance","Scan the SSG QR"),
 ("Sit the WA then the CS","Open book · 1 hour each"),
 ("Submit on the LMS","Upload your answers"),
 ("Sign the record","Assessment Summary Record")],
 kicker="ON ASSESSMENT DAY",color=BLUE,
 synthesis=("REMEMBER","All five steps are mandatory for WSQ funding — missing the digital attendance or the TRAQOM survey can invalidate your claim."))
tile_grid("Criteria for Funding",[
 ("Attendance","A minimum attendance rate of 75%, based on the SSG Digital Attendance record."),
 ("Assessment","Complete both assessment components and be assessed as 'Competent'."),
 ("Digital attendance","Scan the SSG QR code for AM, PM and Assessment on every training day."),
 ("TRAQOM survey","Complete the mandatory TRAQOM course feedback survey on the LMS.")],
 kicker="WSQ FUNDING",cols=2,size=15,accent=AMBER)
tile_grid("The Lab Tools You'll Use",[
 ("Cybersecurity Threat Simulator","alfredang.github.io/cybersecuritysimulator — ten threat modules: phishing, XSS, SQL injection, passwords, social engineering, data leakage."),
 ("Hacklab — Ethical Hacking Simulator","alfredang.github.io/ethnicalhacking — eleven scripted labs across recon, scanning, enumeration and forensics, in a safe simulated terminal."),
 ("FauxBank — Pentest Sandbox","pentest-fauxbank.vercel.app — a fictional bank for guided pentest, a simulated scanner and report generation."),
 ("Cryptography Toolkit","alfredang.github.io/cryptography-toolkit — AES, RSA and ECDSA for the data-protection lab."),
 ("All browser-based","Zero install, no real network traffic, fictional data only — safe to run on the training network."),
 ("Never against real systems","These techniques are for authorised testing only. Never apply them to a system you do not have written authorisation to test.")],
 kicker="LAB ENVIRONMENT · ALL BROWSER-BASED",cols=2,size=14,accent=TEAL)
gallery_slide("The Lab Tools — What You'll See",[
 ("tool-threatsim.png","Cybersecurity Threat Simulator","Ten threat modules: phishing, SQL injection, XSS, passwords, social engineering and data leakage. Used in Lab 2 (mapping the AI threat landscape)."),
 ("tool-fauxbank.png","FauxBank Pentest Sandbox","A fictional bank for guided pentest — IDOR, injection and broken access control — with a simulated scanner and report generation. Used in Lab 10 (pre-deployment red team)."),
 ("tool-hacklab.png","Hacklab Ethical Hacking Simulator","Eleven scripted labs across recon, scanning, enumeration and forensics in a safe simulated terminal. Used in Lab 10 (reconnaissance stage)."),
 ("tool-crypto.png","Cryptography Toolkit","AES, RSA and ECDSA — encryption, key generation and digital signatures. Used in Lab 9 (governing the training data).")],
 kicker="ALL BROWSER-BASED · FICTIONAL DATA ONLY",accent=TEAL,
 note="No installation and no real network traffic. These techniques are for authorised testing only — never apply them to a system you do not have written authorisation to test.")
tile_grid("The Running Case Study — NovaBank",[
 ("One organisation, sixteen labs","Every lab builds on the last for the same fictional Singapore retail bank, so your outputs compose into one governance programme."),
 ("Day 1 · Discover and frame","You inventory the AI estate, map the threats, analyse an agent design and build the business case."),
 ("Day 1 · Build the framework","You run a NIST gap assessment, draft the policy set, design the operating model and apply the PDPA."),
 ("Day 2 · Control the lifecycle","You govern the data, red-team the system, design deployment gates and monitoring, and run an incident."),
 ("Day 2 · Secure the agents","You design the agentic control stack, apply the CSA agentic profile, assess risk and build the roadmap."),
 ("You leave with artefacts","An inventory, a policy set, a RACI, a risk assessment and a 12-month roadmap you can adapt to your own organisation.")],
 kicker="CASE STUDY · NOVABANK",cols=2,size=14,accent=VIOLET)

# ---------------- CORE CONCEPTS ----------------
section("CORE CONCEPTS","Why AI Changes the Security and Governance Picture","")
big_statement("Security governance asks 'is the system protected?'",
 "AI security governance also asks: where did the data come from, what did the model learn, who decided, can we explain it — and what is the agent allowed to do on its own?",
 "THE SHIFT",color=BLUE)
tile_grid("What Is AI Security Governance?",[
 ("A system, not a document","The policies, roles, controls and evidence that keep AI systems secure, lawful and accountable across their whole life."),
 ("Built on what you have","It extends information-security and data-protection governance — it does not replace them, and building it in parallel is a known failure mode."),
 ("Three questions it must answer","What AI do we run? Who is accountable for each system? What happens when one of them behaves badly?"),
 ("Evidence is the output","Governance that cannot produce evidence to a regulator, an auditor or a customer is not governance — it is intention.")],
 kicker="LO1 · DEFINITION",cols=2,size=15)
compare_table("Traditional IT Governance vs AI Security Governance",
 ["Dimension","Traditional IT system","AI system"],
 [["Behaviour","Deterministic — same input, same output","Probabilistic — the same input can produce different output"],
  ["Attack surface","Code, config, infrastructure, credentials","Plus training data, model, prompts, context and tools"],
  ["Change control","Releases go through CAB","Model, prompt and retrieval changes often bypass CAB entirely"],
  ["Failure mode","It crashes or returns an error","It answers confidently and incorrectly"],
  ["Explainability","Read the code","May require dedicated explainability technique"],
  ["Testing","Does it meet the spec?","Plus: is it fair, robust, private, and jailbreak-resistant?"],
  ["Autonomy","None — it does what it is told","An agent decides what to do next, and acts"]],
 kicker="LO1 · WHY THE OLD MODEL DOES NOT FIT",accent=BLUE,
 note="Every row on the right is a control gap in an organisation that has only classical IT governance — which is most organisations.")
tile_grid("The Five New Assets You Must Protect",[
 ("1 · Training and fine-tuning data","Poisoned or unlawfully sourced data becomes permanent model behaviour. Provenance matters before training, not after."),
 ("2 · The model itself","Weights are intellectual property and can be stolen, extracted or inverted to reveal training data."),
 ("3 · Prompts and context","The system prompt is a security control. Whatever enters the context window can influence behaviour."),
 ("4 · The inference endpoint","An unauthenticated or over-permissive endpoint is an open door with a very large blast radius."),
 ("5 · The tools an agent can call","This is the one most organisations miss. The agent's risk is set by what its tools can do, not by what the model says."),
 ("Plus the humans","Over-reliance on confident output is a control failure, not a user error — design the human check deliberately.")],
 kicker="LO1 · THE EXPANDED ATTACK SURFACE",cols=2,size=14,accent=VIOLET)
compare_table("Traditional Threats and Their AI-Era Variants",
 ["Traditional threat","AI-era variant","What changes"],
 [["Phishing","GenAI spear-phishing and deepfakes","Volume and quality: personalised at scale, near-zero cost"],
  ["SQL injection","Prompt injection","Data and instructions share one channel; no parameterisation exists"],
  ["XSS","Unsafe rendering of model output","Markdown links and HTML in output can auto-fetch attacker URLs"],
  ["Weak credentials","Over-scoped agent identity","One shared service account with broad read access to everything"],
  ["Malware","Poisoned package or model artefact","Supply chain now includes models, embeddings and MCP servers"],
  ["Insider data theft","Training-data and context leakage","The model memorises; the context window carries it out"],
  ["Ransomware","Agent with write access at machine speed","No human pace limit on destructive action"]],
 kicker="LO1 · MAPPING WHAT YOU ALREADY KNOW",accent=AMBER,
 note="Your existing security knowledge transfers. What changes is the channel, the speed and the fact that the instruction and the data arrive together.")
tile_grid("The 12 NIST Generative AI Risks",[
 ("CBRN information","Access to materially lowered barriers for weapons-related information."),
 ("Confabulation","Confidently stated false content — the risk most likely to reach your customers."),
 ("Dangerous or violent content","Generation of content that recommends or facilitates harm."),
 ("Data privacy","Leakage or inference of personal data from training data, prompts or outputs."),
 ("Environmental impact","Energy and resource cost of training and serving at scale."),
 ("Harmful bias & homogenisation","Discriminatory outputs; monoculture from one model used everywhere."),
 ("Human-AI configuration","Over-reliance, poor handoff, and humans rubber-stamping model output."),
 ("Information integrity","Synthetic content degrading trust in what is real."),
 ("Information security","Prompt injection, model extraction, expanded attack surface."),
 ("Intellectual property","Training on, or reproducing, protected material."),
 ("Obscene or abusive content","Generation of non-consensual or degrading material."),
 ("Value chain & components","Third-party models, data and packages you did not build and cannot fully inspect.")],
 kicker="LO1 · NIST AI 600-1 GENERATIVE AI PROFILE",cols=3,size=12)
big_statement("The lethal trifecta",
 "Private data access  +  exposure to untrusted content  +  the ability to communicate externally. An agent with all three can be steered by an attacker into leaking data through entirely permitted actions.",
 "LO1 · THE PATTERN BEHIND THE 2025–26 AGENT BREACHES",color=AMBER)
process_map("How a Prompt-Injection Data Leak Actually Happens",[
 ("Attacker plants text","In an e-mail, a web page, a document or a shared file"),
 ("Agent ingests it","Retrieval pulls the content into the context window as ordinary data"),
 ("Model reads it as instruction","There is no separation between data and instructions in a prompt"),
 ("Agent acts","It calls a permitted tool — send e-mail, fetch a URL, render a link"),
 ("Data leaves","Through an authorised channel, with no exploit and no alert")],
 kicker="LO1 · THE ECHOLEAK PATTERN",color=AMBER,
 synthesis=("WHY FILTERS DON'T SAVE YOU","Every step uses permitted functionality. Nothing is 'hacked'. This is why the fix is architectural — break a leg of the trifecta — rather than adding another output filter."))
chart_slide("Where AI Risk Concentrates in a Typical Business",
 ["Shadow AI / unapproved tools","Customer-facing GenAI","Agentic systems with tools",
  "Third-party AI in SaaS","Internally built models"],
 [("Share of organisations reporting incidents (%)",[45,22,18,12,8])],
 kicker="LO1 · WHERE TO LOOK FIRST",accent=BLUE,kind="bar",
 insight="Shadow AI dominates because it is invisible to every control you already run. Illustrative teaching figures — the pattern, not the precise number, is the point: your first control is an inventory, not a firewall.",
 number_format='0"%"')
tile_grid("Business Impact — Speaking the Board's Language",[
 ("Regulatory","PDPA enforcement, MAS supervisory action, and the cost of proving compliance you cannot currently evidence."),
 ("Contractual","Client contracts increasingly require AI assurance. No assurance, no renewal."),
 ("Intellectual property","Confidential strategy, code and customer data pasted into tools that retain and train on it."),
 ("Operational","An agent that acts wrongly at machine speed causes damage faster than a human can notice."),
 ("Decision quality","Confident, wrong output relied upon in credit, hiring or fraud decisions about real people."),
 ("Reputational","'Bank's AI leaked customer data' is a headline that outlives the technical fix.")],
 kicker="LO1 · WHY THIS IS A BOARD ISSUE",cols=2,size=14,accent=AMBER)
section("FRAMEWORKS","The Governance Frameworks You Will Apply","")
process_map("NIST AI RMF 1.0 — The Four Functions",[
 ("GOVERN","Policies, roles, culture, third-party risk — cross-cutting"),
 ("MAP","Frame the context: purpose, actors, risk tolerance"),
 ("MEASURE","Select metrics, test, evaluate, review independently"),
 ("MANAGE","Prioritise, treat, monitor, respond, document residual risk")],
 kicker="LO2 · THE FRAMEWORK YOU WILL ASSESS AGAINST",color=BLUE,
 synthesis=("GOVERN IS NOT A STAGE","GOVERN is infused through the other three. Compliance and evaluation aspects of GOVERN should be integrated into MAP, MEASURE and MANAGE — not done once at the start."))
tile_grid("The Seven Characteristics of Trustworthy AI",[
 ("Valid and reliable","It does what it claims, repeatably, in the conditions it will actually meet."),
 ("Safe","It does not create conditions endangering life, health, property or the environment."),
 ("Secure and resilient","It withstands adversarial input and unexpected conditions, and recovers."),
 ("Accountable and transparent","Someone is answerable, and information about the system is available to those who need it."),
 ("Explainable and interpretable","The mechanism and the meaning of an output can be conveyed appropriately."),
 ("Privacy-enhanced","Anonymity, confidentiality and control over personal data are safeguarded."),
 ("Fair with harmful bias managed","Systemic, computational and human-cognitive bias are identified and managed."),
 ("They trade off","These characteristics interact — optimising one can degrade another. Managing the trade-off is the governance work.")],
 kicker="LO2 · NIST AI 100-1",cols=2,size=14,accent=TEAL)
tile_grid("GOVERN — The Six Categories",[
 ("GOVERN 1 · Policies & processes","Legal requirements understood; trustworthiness in policy; risk tolerance; AI inventory (1.6); safe decommissioning (1.7)."),
 ("GOVERN 2 · Accountability","Roles and responsibilities documented (2.1); AI risk training (2.2); executive leadership takes responsibility (2.3)."),
 ("GOVERN 3 · Diverse teams","Decisions informed by diverse demographics, disciplines and expertise; human-AI oversight roles defined."),
 ("GOVERN 4 · Risk culture","Critical thinking and safety-first mindset; teams document and communicate impacts; testing and incident sharing enabled."),
 ("GOVERN 5 · Stakeholder engagement","Feedback from those outside the developing team is collected, prioritised and integrated."),
 ("GOVERN 6 · Third-party risk","Policies address third-party AI and data risk (6.1) and contingency for third-party failures (6.2).")],
 kicker="LO2 · THE FUNCTION MOST ORGANISATIONS SCORE WORST ON",cols=2,size=13,accent=BLUE)
tile_grid("Singapore's MGF for Generative AI — Nine Dimensions  (K3)",[
 ("1 · Accountability","Allocate responsibility along the AI development chain so end-users have someone answerable."),
 ("2 · Data","Ensure data quality and trusted sources; give clarity on contentious data — personal data and copyright."),
 ("3 · Trusted development & deployment","Best practice in development and evaluation, plus 'food label'-type transparency on safety measures taken."),
 ("4 · Incident reporting","Structures and processes for monitoring, timely notification and remediation."),
 ("5 · Testing & assurance","Third-party testing for independent verification, and common standards so results are comparable."),
 ("6 · Security","GenAI introduces new threat vectors through the models themselves; adapt existing InfoSec frameworks."),
 ("7 · Content provenance","Transparency about where and how content was generated — watermarking and cryptographic provenance."),
 ("8 · Safety & alignment R&D","Invest in alignment with human intention and values; global cooperation across AI safety institutes."),
 ("9 · AI for public good","Democratised access, public-sector adoption, upskilling and sustainability.")],
 kicker="K3 · ETHICAL FRAMEWORKS AND GUIDELINES (SINGAPORE)",cols=3,size=12,accent=VIOLET)
tile_grid("AI Verify — Turning Claims Into Evidence  (A6)",[
 ("Singapore's testing framework","An AI governance testing framework and toolkit from IMDA and the AI Verify Foundation."),
 ("Test, don't assert","It lets an organisation demonstrate through standardised tests that a system behaves as claimed."),
 ("Serves two MGF dimensions","It produces exactly the evidence that 'trusted development and deployment' and 'testing and assurance' call for."),
 ("Why it matters commercially","Independent verification is becoming a procurement requirement, not a differentiator.")],
 kicker="LO2 · SINGAPORE TESTING",cols=2,size=15,accent=TEAL)
compare_table("The Frameworks Side by Side  (A6, K3)",
 ["Framework","Origin","What it gives you","Use it for"],
 [["NIST AI RMF 1.0","NIST (US), voluntary","GOVERN / MAP / MEASURE / MANAGE structure","Your gap assessment and control structure"],
  ["NIST AI 600-1","NIST GenAI profile","12 named GenAI risks + suggested actions","Risk identification for generative systems"],
  ["MGF for GenAI","IMDA / AI Verify (SG)","Nine dimensions of a trusted ecosystem","Aligning to Singapore expectations"],
  ["AI Verify","IMDA / AI Verify (SG)","Testing framework and toolkit","Producing assurance evidence"],
  ["PDPA + PDPC AI Guidelines","PDPC (SG), law + guidance","Lawful basis, notification, accountability","Anything touching personal data"],
  ["CSA Agentic RMF Profile","Cloud Security Alliance","Agentic extensions: AG-GV/MP/MS/MG","Governing agents specifically"],
  ["OWASP ASI Top 10","OWASP","Agentic threat taxonomy","Threat modelling an agent"]],
 kicker="LO2 · A6 · EVALUATING ETHICAL FRAMEWORKS AND GUIDELINES",accent=BLUE,
 note="These are complementary, not competing. NIST gives structure, Singapore gives local expectation and law, CSA and OWASP give the agentic detail NIST does not yet cover.")
decision_map("Does the PDPA Let You Use This Personal Data for AI?  (K3)",
 "Do you have\nmeaningful consent\nfor this purpose?",
 ("YES — proceed on consent","Ensure the notification was meaningful: the individual can understand what data is used and broadly how the system uses it to decide or recommend."),
 ("NO — is an exception available?","Business Improvement Exception for improving or developing an existing product or service. Research Exception for commercial research with public benefit. Record the basis in writing."),
 kicker="K3 · LEGAL REQUIREMENTS FOR DATA PROTECTION",
 note="Whichever route you take, the Accountability Obligation still applies: you must be able to SHOW how you discharge your obligations — policies, the written basis, and the measures protecting individuals.")
chart_slide("Typical Organisational Maturity by RMF Function",
 ["GOVERN","MAP","MEASURE","MANAGE"],
 [("Typical score (of 3)",[1.4,1.1,0.6,1.0]),("Target",[2.5,2.5,2.5,2.5])],
 kicker="LO2 · WHERE THE GAPS USUALLY ARE",accent=VIOLET,kind="column",
 insight="MEASURE scores lowest almost everywhere: organisations with mature InfoSec still rarely test AI systems for bias, robustness or jailbreak resistance. Indicative pattern for teaching — you will produce NovaBank's real scores in Lab 5.")

# ---------------- PER-TOPIC DEEP-DIVE TEACHING SLIDES ----------------
# Extra substantive teaching beyond the concept tiles, keyed by topic number.
def topic_deep_dive(n):
    if n==1:
        tile_grid("Shadow AI — The Governance Problem You Already Have",[
         ("It is the norm, not the exception","Staff adopt AI faster than any approval process. Assume it is happening and go looking, rather than assuming it is not."),
         ("It bypasses every control","No DLP, no logging, no retention limit, no contract, no assurance — because the tool was never registered."),
         ("Where to find it","Expense claims and department cards, SaaS AI features switched on by default, browser extensions, and an honest anonymous staff survey."),
         ("Embedded AI counts","M365 Copilot, Zoom AI Companion, Gong and ServiceNow Now Assist are AI systems even though nobody 'bought AI'."),
         ("The wrong response","Blanket bans push usage further underground. Provide a sanctioned, easy route and the shadow shrinks."),
         ("The right first control","An inventory with named owners. You cannot risk-assess, monitor or govern what you have not listed.")],
         kicker="LO1 · THE FIRST THING YOU WILL FIND",cols=2,size=14,accent=AMBER)
        compare_table("Reading an Incident: What Actually Failed",
         ["Incident","Surface cause","The governance control that was missing"],
         [["EchoLeak (M365 Copilot, CVSS 9.3)","Injected instruction in an e-mail","Untrusted content reached a context with private data and an outbound path"],
          ["Build-pipeline instruction injection","Malicious system-prompt file","Over-scoped token; agent-authored change merged without review"],
          ["Unauthenticated MCP server (CVSS 9.1)","No auth at the tool boundary","Tool exposed without authentication — no model involvement needed"],
          ["Config override in a coding agent","Repo-defined settings overrode user approval","Project-level config trusted above user consent; hooks ran pre-authentication"],
          ["Dependency compromise","Malicious package version","No version pinning; automation multiplied the blast radius in hours"]],
         kicker="LO1 · SIX PUBLIC INCIDENTS, ONE LESSON",accent=AMBER,
         note="In none of these did the model 'go rogue'. Every one was a failure of authorisation, scope or supply chain — which is to say, a governance failure.")
        tile_grid("Building the Business Case",[
         ("Lead with exposure, not technology","Boards fund risk reduction and revenue protection. They do not fund frameworks."),
         ("Quantify one scenario properly","One credible, costed scenario beats a list of ten vague ones. Show your assumptions so they can be challenged."),
         ("Name the regulatory exposure","PDPA penalties and supervisory consequences are concrete numbers a board already understands."),
         ("Use the enabler argument","Governed organisations deploy faster because approval and evidence are routine rather than bespoke each time."),
         ("Ask for one thing","Budget, headcount or mandate. A paper with three asks gets none of them."),
         ("Prepare the tool question","'Why not just buy a tool?' — because a tool with no inventory, no owner and no policy has nothing to enforce.")],
         kicker="LO1 · GETTING IT FUNDED",cols=2,size=14,accent=TEAL)
    if n==2:
        tile_grid("What a Workable Policy Set Looks Like",[
         ("AI Acceptable Use Policy","For everyone. What staff may and may not do, with which tools and which data. Short enough to be read."),
         ("AI Risk Management Standard","Risk tiering, the assessment required per tier, who signs off, and the review cycle."),
         ("Data for AI Standard","Provenance, lawful basis, classification, minimisation, de-identification and retention."),
         ("Model & Agent Development Standard","Threat modelling, testing, red-teaming, human-oversight design and change control."),
         ("Third-Party AI Standard","Vendor due diligence, model provenance questions, sub-processor disclosure and contractual terms."),
         ("One page beats twenty","A policy nobody reads is worse than none — it creates the illusion of control.")],
         kicker="LO2 · FIVE DOCUMENTS, NOT ONE",cols=2,size=14,accent=BLUE)
        tile_grid("The Five Policy Defects That Make Policies Fail",[
         ("Vague scope","'This policy covers AI' — does it cover Copilot? Personal ChatGPT? An AI feature in a SaaS you already own? Say so."),
         ("No approval turnaround","A pathway with no service level is a pathway staff bypass. Commit to a number of days."),
         ("Prohibitions without alternatives","Ban a tool without providing a sanctioned route and usage moves underground, not away."),
         ("No named owner","A policy with no owner and no review date is already obsolete on the day it is issued."),
         ("Unenforceable clauses","If you cannot detect a breach, the clause is advice. Pair each rule with the control that evidences it."),
         ("The test","Could a member of staff read it and know exactly what to do on Monday morning?")],
         kicker="LO2 · WRITE FOR ENFORCEMENT",cols=2,size=14,accent=AMBER)
        process_map("Risk Tiering Drives Everything Else",[
         ("Two questions","Does it decide about a person? Does it touch personal or confidential data?"),
         ("Assign a tier","High, Medium or Low — recorded in the inventory"),
         ("Scale the controls","Assessment depth, testing, oversight and approver rise with the tier"),
         ("Set the review cycle","High reviewed most often; Low on a light-touch cycle"),
         ("Proportionate governance","Effort follows risk, so the process survives contact with the business")],
         kicker="LO2 · PROPORTIONALITY IS WHAT MAKES IT SUSTAINABLE",color=TEAL,
         synthesis=("THE FAILURE MODE","Treating every AI system the same. Full assurance on a meeting-summariser exhausts the goodwill you need for the credit model."))
        tile_grid("Roles and Risk Ownership",[
         ("Board / Executive","Owns AI risk appetite and is accountable for decisions about AI risk (NIST GOVERN 2.3). Not delegable."),
         ("AI Governance Committee","Decides. If it can only advise, it cannot govern — write decision rights into the terms of reference."),
         ("AI Governance Lead","Runs the programme: inventory, assessments, reporting. Coordinates; does not own every system."),
         ("AI System Owner","A named business person accountable for one system's purpose, risk and behaviour."),
         ("Data Protection Officer","Lawful basis, notification, individual rights and breach determination."),
         ("Security Lead","Threat modelling, testing, monitoring, incident response and the kill switch.")],
         kicker="LO2 · WHO DECIDES WHAT",cols=2,size=14,accent=VIOLET)
        tile_grid("The Single-A Rule and Why RACI Fails Without It",[
         ("R — Responsible","Does the work. There can be several."),
         ("A — Accountable","Answers for the outcome. There must be exactly ONE."),
         ("C — Consulted","Two-way: their input is sought before the decision."),
         ("I — Informed","One-way: told after the decision."),
         ("Two A's means none","When two roles are accountable, each assumes the other has it. This is discovered during an incident, at the worst possible moment."),
         ("The stress test","Walk a real incident through your RACI. Every step must produce a named person.")],
         kicker="LO2 · REMOVING AMBIGUITY",cols=2,size=14,accent=BLUE)
    if n==3:
        process_map("The AI Lifecycle as a Control Surface",[
         ("Plan","Purpose · tier · lawful basis"),
         ("Data","Provenance · minimise · de-identify"),
         ("Develop","Threat model · change control"),
         ("Test","Bias · robustness · red-team"),
         ("Deploy","Gate · approvers · rollback"),
         ("Operate","Monitor · log · respond"),
         ("Retire","Revoke · dispose · preserve")],
         kicker="LO3 · EVERY STAGE HAS ITS OWN FAILURE MODE",color=BLUE,
         synthesis=("WHERE GOVERNANCE USUALLY STOPS","At deploy. Which is precisely where AI risk begins to change, because the data, the model and the usage all drift."))
        tile_grid("Data Governance for AI — Decisions Made Before Training",[
         ("Provenance","Where did every field come from, and were we permitted to collect it for this purpose?"),
         ("Lawful basis","Consent, or a recorded exception. Decide and document it before the data moves, not after the model works."),
         ("Minimisation","Does the model actually need this field to do its job? Most training sets carry fields nobody can justify."),
         ("Classification","Direct identifier, quasi-identifier, sensitive attribute or non-personal — the class drives the treatment."),
         ("De-identification","Remove, hash, generalise or retain. Note that encryption is reversible and so remains personal data to the key holder."),
         ("Retention and deletion","How long the training set lives, and what happens to the model when an individual exercises deletion.")],
         kicker="LO3 · GOVERN THE DATA FIRST",cols=2,size=14,accent=TEAL)
        compare_table("Protecting Data — Which Technique, For What",
         ["Technique","What it gives you","Reversible?","Still personal data?","Use it for"],
         [["Deletion / omission","The field is simply gone","No","No","Anything the model does not need"],
          ["Generalisation","Age band instead of date of birth","No","Reduced risk","Quasi-identifiers with predictive value"],
          ["Hashing","A consistent pseudonym","No (but linkable)","Yes — it is pseudonymisation","Joining records without holding the identifier"],
          ["Encryption (AES)","Confidentiality in storage/transit","Yes, with the key","Yes, to the key holder","Protecting data at rest and in transit"],
          ["Asymmetric (RSA)","Safe key exchange","Yes, with the private key","Depends","Moving keys and data between environments"],
          ["Signing (ECDSA)","Integrity and origin proof","N/A","N/A","Proving the training set was not altered after approval"]],
         kicker="LO3 · THE DISTINCTION THAT MATTERS",accent=TEAL,
         note="The most common mistake: calling a hashed or encrypted dataset 'anonymised'. It is pseudonymised, and PDPA obligations continue to apply.")
        tile_grid("Testing an AI System — Beyond 'Does It Answer Correctly?'",[
         ("Accuracy and reliability","Does it perform in the conditions it will actually meet, not just on the test set?"),
         ("Robustness","How does it behave on edge cases, adversarial input and malformed data?"),
         ("Bias and fairness","Does performance differ across groups? Test it; do not assume the data was neutral."),
         ("Privacy leakage","Can training data or another user's data be elicited from the output?"),
         ("Jailbreak resistance","Can the instruction be overridden by crafted input? Assume yes, and measure how easily."),
         ("Grounding","Are claims traceable to a source, or is the model confabulating fluently?")],
         kicker="LO3 · WHAT PRE-DEPLOYMENT ASSURANCE MEANS",cols=2,size=14,accent=VIOLET)
        tile_grid("Red-Teaming an AI System",[
         ("It is adversarial, not functional","Functional testing asks 'does it work?'. Red-teaming asks 'how do I make it misbehave?'"),
         ("Direct prompt injection","Instructions typed by the user to override the system prompt."),
         ("Indirect injection","Instructions hidden in content the agent retrieves — e-mail, documents, web pages, even white-on-white text."),
         ("System-prompt extraction","Getting the agent to reveal its own instructions, which then makes every other attack easier."),
         ("Unsafe tool invocation","Inducing a permitted tool call with attacker-chosen parameters."),
         ("When to run it","Before production, and after every significant change — a new tool, a new data source, a model version bump.")],
         kicker="LO3 · ADVERSARIAL ASSURANCE",cols=2,size=14,accent=AMBER)
        tile_grid("Monitoring — Because AI Risk Changes After Launch",[
         ("Quality","Accuracy and groundedness drift as the world and the data move away from the training distribution."),
         ("Safety","Refusal rate and injection-detection rate. A metric with no threshold is a dashboard, not a control."),
         ("Operations","Latency, cost per call and error rate — the signals that surface a problem before customers do."),
         ("Security","Permission escalations, out-of-hours access, anomalous tool sequences and outbound sends."),
         ("Every alert needs an owner","Who is notified, what they check first, and the condition under which the system is disabled."),
         ("Time to disable","Measure it. You will be asked during an incident, and 'we're not sure' is not an answer.")],
         kicker="LO3 · POST-DEPLOYMENT CONTROL",cols=2,size=14,accent=BLUE)
        process_map("AI Incident Response",[
         ("Detect","An alert, a customer report, or someone noticing"),
         ("Contain","Disable or restrict — decided in advance, by a named authority"),
         ("Preserve","Prompts, tool calls, model version, config, the injected content"),
         ("Assess","Scope, individuals affected, data classes, notification threshold"),
         ("Notify","PDPC, individuals, regulator, board — per your determination"),
         ("Learn","Root cause, improvement actions with owners and dates")],
         kicker="LO3 · EXTEND THE PLAN YOU ALREADY HAVE",color=AMBER,
         synthesis=("THE QUESTION THAT EXPOSES THE GAP","If your agent started leaking data at 09:00 on a Saturday, who switches it off, and how long does it take? Most organisations cannot answer."))
    if n==4:
        tile_grid("Why Agents Need Their Own Controls",[
         ("Risk follows the tools","The model's words are not the risk. What the agent can DO — send, write, pay, freeze, delete — is the risk."),
         ("No human pace limit","An agent repeats an action hundreds of times in the time a person takes to notice one."),
         ("Memory persists","Poisoned content that enters agent memory keeps influencing behaviour long after the original input."),
         ("Delegation diffuses accountability","When an agent instructs a sub-agent, who approved that instruction? Usually nobody."),
         ("Identity is often shared","One service account for the agent means every action is attributable to 'the app', not to an owner."),
         ("Frameworks lag","The NIST AI RMF treats a read-only recommender and an autonomous executor identically. The CSA profile closes that gap.")],
         kicker="LO4 · THE AGENTIC DIFFERENCE",cols=2,size=14,accent=VIOLET)
        process_map("Defence in Depth for an Agent",[
         ("Identity","Its own identity, scoped short-lived credentials"),
         ("Guardrails","Input and output filtering at the model boundary"),
         ("Permission ladder","Deny / ask / allow, evaluated before the tool runs"),
         ("Sandbox","Restricted filesystem, egress and workspace"),
         ("Human gate","Approval on irreversible and outbound actions"),
         ("Audit trail","Who approved what, when and why — reconstructable")],
         kicker="LO4 · SIX LAYERS, NOT ONE FILTER",color=BLUE,
         synthesis=("PERMISSION IS INFRASTRUCTURE, NOT PROMPT","A rule written in the system prompt is a request. A rule enforced before the tool executes is a control. Only the second survives an injection."))
        compare_table("What Each Layer Can and Cannot Do",
         ["Layer","Stops","Does NOT stop"],
         [["Content guardrails","Toxic text, obvious jailbreaks, PII in output","Tool authorisation, credential scope, multi-step exfiltration"],
          ["Permission ladder","Unauthorised tool calls, dangerous parameters","A permitted tool used for an attacker's purpose"],
          ["Sandbox","Filesystem and network escape, blast radius","Anything the agent is legitimately allowed to reach"],
          ["Human approval","Irreversible action taken silently","Anything below the approval threshold; fatigue-driven rubber-stamping"],
          ["Audit trail","Nothing — it is detective, not preventive","Any attack; it tells you what happened afterwards"],
          ["Architecture (break the trifecta)","Whole classes of exfiltration","Threats inside the remaining capability set"]],
         kicker="LO4 · HONEST ABOUT LIMITS",accent=AMBER,
         note="No single layer is sufficient. The architectural choice — not combining private data, untrusted content and an outbound channel in one agent — does more than all the filters put together.")
        tile_grid("The CSA Agentic Extensions to the NIST AI RMF",[
         ("AG-GV.1 Autonomy tier classification","Four tiers with escalating oversight, because a recommender and an executor are not the same risk."),
         ("AG-GV.2 Delegation accountability","Document oversight boundaries, escalation triggers and the lineage from an action back to a named human."),
         ("AG-GV.3 Agent lifecycle registry","A live inventory of agent authorities, tool access, delegation relationships and review schedule."),
         ("AG-MP.1 / MP.2 Tool and consequence analysis","Classify tools by consequence and reversibility; map which tool sequences lead to which real-world outcomes."),
         ("AG-MS.1 Behavioural telemetry","Action velocity, permission escalation rate, cross-boundary calls, delegation depth, exception rate."),
         ("AG-MG.1 / MG.3 Incident and decommissioning","Playbooks for agent compromise, hijack, runaway and delegation-chain compromise; structured retirement.")],
         kicker="LO4 · CLOSING THE FOUR STRUCTURAL GAPS",cols=2,size=13,accent=TEAL)
        tile_grid("Autonomy Tiering — Tier by Capability, Not by Title",[
         ("Tier 1 · Supervised","Proposes only; a human acts. Lightest oversight — but confirm it truly cannot act."),
         ("Tier 2 · Constrained","Acts within a fixed allowlist of reversible actions. Periodic review of the allowlist."),
         ("Tier 3 · Broad with monitoring","Wide capability with runtime telemetry and thresholds. Continuous monitoring is mandatory, not optional."),
         ("Tier 4 · Full autonomy","Highest obligation: pre-authorised containment, tested kill switch, named standing authority to disable."),
         ("The common error","Tiering by how the agent is described. 'Just a helpdesk assistant' that resets passwords is not Tier 1."),
         ("Tier drives obligation","If the oversight requirement does not rise with the tier, the tiering is decorative.")],
         kicker="LO4 · AG-GV.1",cols=2,size=14,accent=VIOLET)
        tile_grid("Approval Fatigue — When a Control Stops Being One",[
         ("The failure","When almost every action prompts for approval, approvers stop reading and start clicking."),
         ("It gets worse with time","Approval discipline decays as session count rises — the pattern is well documented across agent deployments."),
         ("Measure the edit rate","Not the approval rate. How often does the human actually change or reject the proposed action?"),
         ("Auto-allow the routine","Reserve the gate for irreversible and outbound actions; let low-consequence reads through."),
         ("Watch for denial clusters","Three consecutive denials, or twenty in a session, should escalate to a human review of the agent itself."),
         ("Design the volume deliberately","State the number of approvals per day above which your gate stops being a control.")],
         kicker="LO4 · A CONTROL THAT DEGRADES IF YOU IGNORE IT",cols=2,size=14,accent=AMBER)
        process_map("Structured AI Risk Assessment",[
         ("Scope","The system, version, data, tools, users and boundary"),
         ("Identify","Risks as cause → event → consequence statements"),
         ("Rate inherent","Likelihood x impact, before controls"),
         ("Map controls","Only controls that actually exist or are committed"),
         ("Rate residual","Honestly — filters cut likelihood, rarely impact"),
         ("Treat and own","Treat / tolerate / transfer / terminate, with an owner and a date")],
         kicker="LO4 · THE DOCUMENT A REGULATOR WILL ASK FOR",color=BLUE,
         synthesis=("THE MOST COMMON DEFECT","Risks written as single words. 'Prompt injection' is a cause. A risk statement names the cause, the event and the business consequence."))
        chart_slide("Sequencing the Roadmap by Risk Reduction per Unit of Effort",
         ["AI inventory + owners","Acceptable use policy","Logging enabled","Risk assess High-tier",
          "Deployment gates","Agentic control stack","Independent assurance"],
         [("Risk reduction",[9,8,7,8,7,6,4]),("Effort",[2,2,3,5,5,7,8])],
         kicker="LO4 · WHY PHASE 1 IS CHEAP AND HIGH-VALUE",accent=TEAL,kind="column",
         insight="The first three initiatives are low effort and high risk reduction — do them first. Independent assurance matters, but it cannot precede having something to assure. Relative teaching values, not effort estimates for your organisation.")
        tile_grid("Metrics That Keep a Programme Alive",[
         ("Inventory coverage","Percentage of known AI systems registered with a named owner. Your foundational metric."),
         ("Assessment coverage","Percentage of High-risk systems with a current risk assessment."),
         ("Policy attestation","Percentage of staff who have attested to the acceptable-use policy."),
         ("Time to disable","Minutes to disable a live agent — tested, not estimated."),
         ("Incidents and MTTR","Count of AI incidents and mean time to respond."),
         ("Report them upward","Governance that is not measured and reported does not survive the next budget cycle.")],
         kicker="LO4 · PROVING IT WORKS",cols=2,size=14,accent=BLUE)


# ---------------- RESPONSIBLE-AI FOUNDATIONS (TSC K1-K5 / A2, A4, A8, A9) ----------------
# The approved TSC for this course (ICT-INT-0055-1.1) assesses responsible-AI knowledge and
# abilities — environmental impact, bias and societal impact, the privacy-performance
# trade-off. Security governance is the vehicle; these are the assessed outcomes, so they are
# taught explicitly here.
section("RESPONSIBLE AI","Ethics, Bias, Privacy and Sustainability in AI Governance","")
big_statement("Security governance and responsible AI are the same programme.",
 "A system can be perfectly secure and still be unfair, opaque, privacy-invasive or wasteful. The governance framework you build has to carry both — which is why this course assesses both.",
 "WHY THIS SECTION",color=VIOLET)
tile_grid("Ethical Considerations Across AI Development  (K5)",[
 ("Bias","Consider it at every stage — data provenance and representation, the fairness measure chosen at design, group-wise testing before release, and drift monitoring after it."),
 ("Privacy","Minimisation and purpose limitation before training; lawful basis under the PDPA; de-identification; access control; retention limits; and what happens on a deletion request."),
 ("Transparency","Disclose that AI is in use, what it does and its limits. Meaningful notification where personal data is used; explanation where a decision affects someone."),
 ("Accountability","A named human accountable for each system — NIST GOVERN 2.1 and 2.3, and the first MGF dimension. In practice: an inventory with owners and one Accountable per control."),
 ("They trade off","Stronger privacy usually costs accuracy; more explainability constrains model choice. Managing and RECORDING the trade-off is the governance work."),
 ("Proportionality","Risk-tier each system and scale assessment, testing and oversight to the tier, so effort follows risk and the process survives contact with the business.")],
 kicker="K5 · ETHICS IN AI DEVELOPMENT",cols=2,size=14,accent=VIOLET)
process_map("How Bias Enters an AI System  (K2, A4)",[
 ("Historical data","Past decisions encoded as ground truth"),
 ("Representation","Under-represented groups modelled worse"),
 ("Proxy features","Postal code standing in for age or income"),
 ("Measurement","Some groups' data captured less accurately"),
 ("Feedback loop","Its own outputs become tomorrow's training data")],
 kicker="K2 · ANALYSING BIAS AND FAIRNESS",color=AMBER,
 synthesis=("THE POINT","Bias is rarely a biased developer. It is ordinary data and ordinary design choices reproducing an existing pattern — which is why it must be TESTED for, not assumed absent."))
tile_grid("Implications of Biased AI Algorithms  (K2, A8)",[
 ("On the individual","A concrete unfair outcome — a delayed appointment, a declined loan, a screened-out application — usually falling on the person least able to challenge it."),
 ("On groups and cultures","One model decides every case identically, so an error is applied SYSTEMATICALLY to a whole group rather than randomly. Scale turns small bias into population-level harm."),
 ("On minority groups","Smaller groups are under-represented in training data, so error rates are higher exactly where the ability to contest a decision is often lowest."),
 ("On language and culture","Systems perform best in the language and cultural context they were trained on — a live concern in multilingual Singapore."),
 ("On the organisation","Regulatory exposure, contractual and legal risk, loss of trust, and remediation cost far above the cost of testing beforehand."),
 ("Compounding","Feedback loops and algorithmic monoculture — many organisations using the same model — entrench the same bias across a whole market.")],
 kicker="K2 / A8 · SOCIETAL AND CULTURAL IMPACT",cols=2,size=13,accent=AMBER)
compare_table("Fairness Is Not One Thing — Choose the Measure  (A4)",
 ["Fairness measure","What it requires","Use it when","The catch"],
 [["Demographic parity","Equal positive rates across groups","Outreach, access to opportunity","Ignores genuine differences in need"],
  ["Equal opportunity","Equal true-positive rates across groups","Screening where a miss is the harm","Needs reliable ground-truth labels"],
  ["Equalised odds","Equal true- AND false-positive rates","High-stakes decisions about people","The hardest to satisfy in practice"],
  ["Calibration","A given score means the same for every group","Risk scoring and pricing","Can coexist with unequal error rates"],
  ["Individual fairness","Similar people treated similarly","Case-by-case review","Requires defining 'similar'"]],
 kicker="A4 · ASSESSING BIAS AND FAIRNESS",accent=AMBER,
 note="These measures are mathematically incompatible — you cannot satisfy all of them at once. Choosing which one fits the use case, and recording why, is a governance decision, not a technical one.")
process_map("Mitigating Bias Across the Lifecycle  (A5)",[
 ("Examine the data","Provenance, representation, proxy features"),
 ("Choose the measure","Fairness metric fixed BEFORE testing"),
 ("Test by group","Aggregate accuracy hides group error"),
 ("Human review","Where a decision materially affects a person"),
 ("Monitor for drift","Fair at launch is not fair forever")],
 kicker="A5 · IMPLEMENTING MITIGATION",color=TEAL,
 synthesis=("VERIFY, DON'T ASSERT","A mitigation is only real if you can show the differential error rate fell. Use AI Verify to produce that evidence, and have it reviewed by someone who did not build the model (NIST MEASURE 1.3)."))
compare_table("The Privacy–Performance Trade-off  (A2)",
 ["Technique","Privacy gain","Performance cost","When to use it"],
 [["Drop unneeded fields","High — the data is simply gone","Usually none","Always. Do this first; it is close to free"],
  ["Generalise (bands, districts)","Moderate to high","Small to moderate","Quasi-identifiers that still carry signal"],
  ["Hash / pseudonymise","Limited — still personal data","None for modelling","Joining records without holding the identifier"],
  ["Aggregation","High","Moderate","Reporting and analytics, not per-person decisions"],
  ["Differential privacy","Very high, and measurable","Material, grows with the guarantee","Secondary analysis; rarely a safety-critical model"],
  ["Federated learning","High — data stays local","Moderate, plus complexity","Data that legally cannot be centralised"]],
 kicker="A2 · EVALUATING THE TRADE-OFF",accent=BLUE,
 note="There is no option that is free on both sides. Decide where on the curve you sit, record the reasoning, and have the data owner and the DPO both sign it — an unrecorded trade-off is a decision nobody made.")
chart_slide("Where the Energy Goes  (K1, A9)",
 ["Training (one-off)","Inference year 1","Inference year 2","Inference year 3"],
 [("Cumulative energy share (%)",[100,145,190,235])],
 kicker="K1 · ENVIRONMENTAL IMPACT OF GENERATIVE AI",accent=TEAL,kind="column",
 insight="Training is the visible one-off cost; inference is small per request but recurs on every request, so for a widely used deployed system it overtakes training and keeps growing. Illustrative shape for teaching the principle — measure your own.",
 number_format='0"%"')
tile_grid("Assessing the Energy Footprint  (K1, A9)",[
 ("What to measure","Energy per inference (kWh per 1,000 requests) x request volume, plus training or fine-tuning as a periodic cost."),
 ("Beyond electricity","Water for data-centre cooling, and the embodied carbon of manufacturing and disposing of the hardware."),
 ("Where and when it runs","Grid carbon intensity varies by region and by time of day — the same workload has different footprints in different places."),
 ("Right-size the model","The single biggest lever. A large reasoning model used for routine classification multiplies energy for no benefit."),
 ("Cut needless inference","Cache repeated answers, batch, shorten prompts and outputs, and remove AI where a rule or a lookup would do."),
 ("Measure and report","NIST AI 600-1 lists Environmental Impacts as a GenAI risk — it belongs in the risk register, not only the sustainability report.")],
 kicker="K1 / A9 · ESTIMATING THE FOOTPRINT",cols=2,size=13,accent=TEAL)
tile_grid("Communicating Capabilities and Limitations Honestly  (K4)",[
 ("Why it matters","A model sounds equally confident whether it is right or wrong. Users calibrate trust on how you describe the system, so overstating it directly causes over-reliance."),
 ("Confabulation","NIST's name for confidently stated false content. It is a property of how these models work, not a bug you can fully remove."),
 ("Human-AI configuration","The NIST risk covering over-reliance, poor handoff and humans rubber-stamping output. Design the human check; do not assume it."),
 ("Say what it is NOT for","The most useful disclosure. 'This does not provide medical, legal or financial advice' prevents more harm than a capability list."),
 ("Publish a model or system card","Purpose, data, performance, limitations and intended use — the 'food label' the MGF's trusted-development dimension calls for."),
 ("Give an escalation route","An obvious, easy path to a human. Transparency without a route to a person is only a disclaimer.")],
 kicker="K4 · TRUSTWORTHY AND TRANSPARENT COMMUNICATION",cols=2,size=13,accent=VIOLET)
process_map("Integrating Ethics Into Development and Evaluation  (A1)",[
 ("Inventory and own","Every AI system registered, with a named accountable human"),
 ("Risk-tier","Decides about a person? Touches personal data?"),
 ("Embed at each stage","Data, build, test, deploy, operate, retire"),
 ("Evaluate before AND after","Pre-deployment testing; post-deployment monitoring"),
 ("Measure the programme","Coverage, assessments, time to disable, incidents")],
 kicker="A1 · DEVELOPING THE STRATEGY",color=BLUE,
 synthesis=("THE STRATEGY IN ONE LINE","Ethics is integrated when it is a GATE and a METRIC — something a system must pass to ship and something you report on afterwards — not a principle in a document nobody consults."))
tile_grid("Applying Ethical Guidelines to a Live System  (A3)",[
 ("Fairness","Test performance across the groups the system affects, and give anyone it serves poorly an accessible route to a human."),
 ("Accountability","A distinct identity per system with scoped credentials, a named owner in the inventory, and logging that makes actions attributable."),
 ("Transparency","Disclose AI use, state what the system is not for, publish a model or system card, and notify meaningfully where personal data is used."),
 ("Sustainability","Right-size the model, cut needless inference, and measure and report the footprint."),
 ("Security","Break the lethal trifecta, apply least privilege and a permission ladder, and gate irreversible actions."),
 ("Proportionality","Scale all of the above to the system's risk tier — full assurance on a meeting summariser exhausts the goodwill you need for the credit model.")],
 kicker="A3 · APPLYING THE GUIDELINES",cols=2,size=13,accent=TEAL)
tile_grid("Championing Responsible AI in Your Organisation  (A7)",[
 ("Role-specific training","Generic awareness changes nothing. Give each group its real tasks — what an administrator may enter, when a professional must not rely on output."),
 ("Make the right path the easy path","A sanctioned toolset, a fast approval route with a real service level, and templates. Governance that obstructs is governance that is bypassed."),
 ("Use your own near-misses","Your organisation's own incident is far more persuasive than any external example."),
 ("Label AI-generated content","Content provenance — MGF dimension 7. Disclose, and review before publication where the audience could be harmed."),
 ("Report results, not intentions","Commit publicly to testing before deployment and monitoring after it, and publish what the tests found."),
 ("Lead with the enabler argument","Governed organisations deploy FASTER, because approval and evidence become routine instead of bespoke every time.")],
 kicker="A7 · CHAMPIONING BEST PRACTICE",cols=2,size=13,accent=BLUE)


LAB_STAGES = {1: [('Read the three sources', 'IT · procurement · survey'), ('Flag every AI system', 'description, not the flag'), ('Merge into one register', 'owner · data · model'), ('Tier autonomy and risk', 'Tier 1–4 · High/Med/Low'), ('Report the shadow AI', 'the systems nobody owns')], 2: [('Run the simulator modules', 'phishing · SQLi · XSS'), ('Record what you observed', 'the mechanism, not the term'), ('Test passwords and social', 'entropy · crack time'), ('Model the data leakage', 'toggle each control'), ('Map to AI-era variants', 'and the business impact')], 3: [('Read the design brief', 'every source, every action'), ('Test the three legs', 'private · untrusted · outbound'), ('Write the attack path', 'the EchoLeak pattern'), ('Map five NIST risks', 'with business impact'), ('Break a leg, not a filter', 'then state residual risk')], 4: [('Open the cost pack', 'four cost categories'), ('Count your High-risk systems', 'from the Lab 1 inventory'), ('Cost one real scenario', 'likelihood × impact'), ('Pick three investments', 'cost · time · risk cut'), ('Write the board paper', 'one page, one ask')], 5: [('Read the evidence pack', 'policies and interviews'), ('Score GOVERN', '0–3, cite the evidence'), ('Score MAP and MEASURE', 'expect MEASURE lowest'), ('Score MANAGE', 'residual risk accepted?'), ('Build the gap register', 'risk cut per unit effort')], 6: [('Set the scope precisely', 'which tools are covered?'), ('Write permitted use', 'data classes and logging'), ('Write prohibited use', 'six concrete, detectable rules'), ('Define the approval path', 'with a real turnaround'), ('Map to MGF dimensions', 'and peer-review it')], 7: [('Reuse existing forums', 'before creating a new one'), ('Write decision rights', 'advice is not governance'), ('Define six roles', 'one owned decision each'), ('Build the RACI', 'exactly one A per control'), ('Stress-test on an incident', 'a name at every step')], 8: [('Read the four scenarios', 'what happens to a person?'), ('Find the lawful basis', 'consent or an exception'), ('Determine notification', 'meaningful, not formal'), ('Apply accountability', 'what can you show?'), ('Write the determination', 'and the risk to escalate')], 9: [('Classify all 24 fields', 'identifier · quasi · sensitive'), ('Minimise ruthlessly', 'does churn need it?'), ('Choose the method per field', 'remove · hash · generalise'), ('Test AES, RSA and ECDSA', 'encryption is reversible'), ('Set retention and deletion', 'and record the basis')], 10: [('Read the authorisation', 'sandbox, fictional data only'), ('Work the guided pentest', 'evidence as you go'), ('Compare with the scanner', 'each misses what the other finds'), ('Add the AI test cases', 'injection · extraction · leakage'), ('Report with conditions', 'a verdict is not assurance')], 11: [('List the gate evidence', 'documents, not intentions'), ('Assign approvers', 'reuse the Lab 7 RACI'), ('Find the log anomalies', 'burst · refusals · injection'), ('Set ten thresholds', 'a metric needs a number'), ('Define the response', 'who acts, and when to disable')], 12: [('Take your role', 'and answer only for it'), ('Contain the incident', 'who authorises, how fast?'), ('Preserve the evidence', 'before anything is reset'), ('Assess the data impact', 'and the PDPA threshold'), ('Notify, then learn', 'five actions with owners')], 13: [('Classify all twelve tools', 'consequence and reversibility'), ('Give the agent an identity', 'scoped, short-lived'), ('Set deny / ask / allow', 'irreversible ⇒ ask or deny'), ('Bound the blast radius', 'filesystem · egress · creds'), ('Gate the human approvals', 'and watch for fatigue')], 14: [('Read the agent estate', 'five agents, real capabilities'), ('Tier by capability', 'not by job title'), ('Document the delegation', 'lineage to a named human'), ('Build the agent registry', 'no blank kill-switch'), ('Set telemetry and playbooks', 'baseline then threshold')], 15: [('Scope it tightly', 'system, version, tools, users'), ('Write cause → event → consequence', 'not a single word'), ('Rate inherent risk', 'justify every 4 and 5'), ('Map only real controls', 'then rate residual honestly'), ('Treat, own and date it', 'High ⇒ treat or accept')], 16: [('Consolidate every artefact', 'gaps and rated risks'), ('Score maturity today', 'evidence per score'), ('Phase 1 — foundation', 'inventory · policy · logging'), ('Phase 2 and 3', 'controls, then assurance'), ('Metrics and the one ask', 'baseline and target')]}

# ---------------- TOPICS + ACTIVITIES ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    cons=t["concepts"]
    half=(len(cons)+1)//2
    for ci,chunk in enumerate([cons[:half],cons[half:]]):
        if not chunk: continue
        suffix="" if ci==0 else " (continued)"
        tile_grid(f"Key Concepts — {t['title']}{suffix}", chunk,
                  kicker=f"TOPIC WEIGHTING {t['weighting']}", cols=2, size=14)
    topic_deep_dive(t["num"])
    acts=TOPIC_ACTS[t["num"]]
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        _lbl=("—" if not g else (f"Lab {g[0]['num']}" if g[0]['num']==g[-1]['num']
                                 else f"Labs {g[0]['num']}–{g[-1]['num']}"))
        cards.append((CARD_COLORS[gi], _lbl,
                      [a["title"] for a in g] if g else ["—"]))
    cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"],
                          kicker=f"TOPIC {t['code']} · HANDS-ON",
                          objective=a.get("objective"), test=a.get("test"))
        # Process-map stages are AUTHORED short labels (LAB_STAGES), not truncated
        # sentences — a chip must read completely, never trail off in an ellipsis.
        stg=list(LAB_STAGES.get(a["num"],[]))
        if len(stg)>=3:
            process_map(f"How Lab {a['num']} Runs", stg,
                        kicker=f"LAB {a['num']} · PROCESS MAP", color=TEAL,
                        synthesis=("THE POINT", a["objective"]))
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("LO1 · Analysed foundations and risks","You inventoried an AI estate including shadow AI, mapped traditional threats to their AI-era variants, analysed a live agent design against the lethal trifecta and the NIST GenAI risks, and built a costed business case."),
 ("LO2 · Evaluated frameworks, built the structure","You ran a NIST AI RMF gap assessment on real evidence, drafted an enforceable policy set mapped to Singapore's MGF, designed a governance operating model with a single-A RACI, and applied the PDPA and PDPC AI guidelines to four scenarios."),
 ("LO3 · Developed lifecycle controls","You governed a training dataset end to end, red-teamed a system before deployment, designed deployment gates and a monitoring regime with thresholds, and ran an AI incident from detection to notification."),
 ("LO4 · Implemented agentic controls and a roadmap","You designed a five-layer agentic control stack, applied the CSA agentic extensions, completed a twelve-risk assessment, and produced a 12-month roadmap with metrics.")],
 kicker="LEARNING OUTCOMES",cols=2,size=13)
tile_grid("The Artefacts You Are Taking Away",[
 ("AI system inventory","A template and a worked example you can populate for your own organisation on Monday."),
 ("Policy set","An acceptable-use policy and a risk-management standard, mapped to MGF dimensions and RMF subcategories."),
 ("Operating model and RACI","A committee terms of reference and a twelve-control RACI that passes the single-A rule."),
 ("Risk assessment","A twelve-risk assessment with a 5x5 matrix, residual ratings, treatments and owners."),
 ("Agentic control stack","A layered specification with a deny/ask/allow tool policy and an agent registry."),
 ("12-month roadmap","Three phases, sequenced by dependency, with six measurable programme metrics.")],
 kicker="WHAT YOU CAN REUSE",cols=2,size=14,accent=TEAL)
tile_grid("Your First 30 Days Back at Work",[
 ("Week 1 · Find out what you run","Start the inventory. Expense claims, SaaS admin consoles and an honest staff survey will find more than IT records will."),
 ("Week 1 · Name owners","Every AI system gets a named accountable human. Systems with no owner are your first finding."),
 ("Week 2 · Issue interim guidance","One page: what staff may and may not put into AI tools, and the route to ask. Do not wait for the full policy set."),
 ("Week 3 · Turn on logging","You cannot investigate what you did not record. Logging is cheap and it is the precondition for everything else."),
 ("Week 4 · Assess your highest-risk system","One proper risk assessment beats a portfolio-wide survey. Use the Lab 15 template."),
 ("Then · Take it to the board","Use your Lab 4 business case and Lab 16 roadmap. Ask for one thing.")],
 kicker="MAKE IT REAL",cols=2,size=14,accent=AMBER)
tile_grid("Continue Your Learning",[
 ("NIST AI RMF","nist.gov/itl/ai-risk-management-framework — the framework, the Playbook and the GenAI profile (AI 600-1)."),
 ("AI Verify Foundation","aiverifyfoundation.sg — the Model AI Governance Framework for GenAI and the testing toolkit."),
 ("PDPC","pdpc.gov.sg — the Advisory Guidelines on the Use of Personal Data in AI Recommendation and Decision Systems."),
 ("OWASP","The Agentic Security Initiative Top 10 — the threat taxonomy for agents."),
 ("Cloud Security Alliance","The Agentic AI NIST RMF Profile — AG-GV, AG-MP, AG-MS and AG-MG control extensions."),
 ("Practise the labs again","Re-run the labs against your own organisation's real AI estate — that is where the value is realised.")],
 kicker="NEXT STEPS",cols=2,size=14)
tile_grid("Recommended Courses",[(rc,"") for rc in C.RECOMMENDED_COURSES],
 kicker="CONTINUE WITH TERTIARY INFOTECH",cols=1,size=15)
content("Support",[
 "If you have any enquiries during or after the class, you can contact us below.",
 "Email: enquiry@tertiaryinfotech.com","Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg","LMS / TMS: https://lms-tms.tertiaryinfotech.com"],kicker="WE'RE HERE TO HELP")
tile_grid("Assessment",[
 ("Written Assessment (SAQ)","1 hour · open book · short-answer knowledge questions on AI security governance."),
 ("Case Study (CS)","1 hour · open book · one coherent governance case study drawn from the labs."),
 ("Digital attendance","Remember to take the Assessment digital attendance (TRAQOM) before you start."),
 ("Submit on the LMS","Upload your completed answers at https://lms-tms.tertiaryinfotech.com/")],
 kicker="WRAP-UP",cols=2,size=15)
process_map("Assessment Flow",[
 ("TRAQOM survey","Scan the QR on the LMS"),
 ("Digital attendance","Scan the SSG QR"),
 ("Sit the WA then the CS","Open book · 1 hour each"),
 ("Submit on the LMS","Upload your answers"),
 ("Sign the record","Assessment Summary Record")],
 kicker="ON ASSESSMENT DAY",color=BLUE,
 synthesis=("REMEMBER","All five steps are mandatory for WSQ funding — missing the digital attendance or the TRAQOM survey can invalidate your claim."))
tile_grid("Digital Attendance (Mandatory)",[
 ("Three times a day","Take the AM, PM and Assessment digital attendance — mandatory for every WSQ-funded course."),
 ("Trainer shows the QR","The trainer or administrator displays the digital attendance QR code from the SSG portal."),
 ("Scan and submit","Scan the QR code with your mobile phone camera and submit your attendance."),
 ("75% minimum","A minimum of 75% attendance is required to be eligible for assessment and funding.")],
 kicker="TRAQOM · SSG DIGITAL ATTENDANCE",cols=2,size=15)
big_statement("Thank You!","You can now analyse AI business risk, build a governance framework, apply controls across the AI lifecycle, secure agentic AI and deliver an implementation roadmap.","GOVERN WELL",color=TEAL)

# ---------------- restrained motion pass ----------------
DIVIDER_IDS=set()
for s in prs.slides:
    txts=[sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
    joined=" ".join(txts)
    is_div=any(k in joined for k in ("COURSE ADMINISTRATION","CORE CONCEPTS","FRAMEWORKS","WRAP-UP")) or \
           any(t["title"] in joined and f"TOPIC {t['code']}" in joined for t in C.TOPICS)
    _transition(s,"push" if is_div else "fade",speed="fast" if not is_div else "med")

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
